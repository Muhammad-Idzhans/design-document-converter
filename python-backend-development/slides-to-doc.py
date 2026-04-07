import os
import sys
import json
import time
import base64
import requests
import shutil
import tempfile
from pathlib import Path
from pptx import Presentation
from dotenv import load_dotenv
from docx import Document
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
try:
    from openai import AzureOpenAI
except ImportError:
    print("Error: The 'openai' library is required to use Microsoft Foundry Native Agents.")
    print("Please run: pip install openai")
    sys.exit(1)

# Removed strict win32com import requirement to allow Linux cross-platform execution

# ──────────────────────────────────────────────
# Configuration
# ──────────────────────────────────────────────
load_dotenv()

CONTENT_UNDERSTANDING_ENDPOINT = os.getenv("CONTENT_UNDERSTANDING_ENDPOINT", "").rstrip("/")
CONTENT_UNDERSTANDING_KEY = os.getenv("CONTENT_UNDERSTANDING_KEY", "")
API_VERSION = "2025-11-01"

# Azure OpenAI Credentials for Vision Analysis
AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT", "").rstrip("/")
AZURE_OPENAI_KEY = os.getenv("AZURE_OPENAI_KEY", "")
AZURE_OPENAI_DEPLOYMENT_NAME = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME", "gpt-4o")

# Azure OpenAI Credentials for Agents (Initial Project)
AGENT_OPENAI_ENDPOINT = os.getenv("AGENT_OPENAI_ENDPOINT", "").rstrip("/")
AGENT_OPENAI_KEY = os.getenv("AGENT_OPENAI_KEY", "")
ORCHESTRATOR_DEPLOYMENT = os.getenv("ORCHESTRATOR_DEPLOYMENT", "gpt-4.1")
WRITER_DEPLOYMENT = os.getenv("WRITER_DEPLOYMENT", "gpt-4.1")
AGENT_ASSISTANT_ID = os.getenv("AGENT_ASSISTANT_ID", "") # Master Agent from Microsoft Foundry UI

SCRIPT_DIR = Path(__file__).parent
PROJECT_DIR = SCRIPT_DIR.parent
OUTPUT_DIR = PROJECT_DIR / "extraction-output"


# ──────────────────────────────────────────────
# Part 1.5: PPTX to PDF Conversion
# ──────────────────────────────────────────────
def convert_pptx_to_pdf(pptx_path: str, pdf_path: str):
    """
    Convert PPTX to PDF using LibreOffice (Linux/Universal) with a graceful fallback to MS PowerPoint COM.
    """
    import os, subprocess
    print(f"\n[CONVERT] Converting PPTX to PDF for Azure Analysis (Cross-Platform)...")
    abs_pptx = os.path.abspath(pptx_path)
    abs_pdf = os.path.abspath(pdf_path)
    
    # Try LibreOffice universally first
    try:
        if os.name == 'posix':
            subprocess.run(['soffice', '--headless', '--convert-to', 'pdf', '--outdir', os.path.dirname(abs_pdf), abs_pptx], check=True)
            expected_out = os.path.join(os.path.dirname(abs_pdf), os.path.splitext(os.path.basename(abs_pptx))[0] + ".pdf")
            if expected_out != abs_pdf and os.path.exists(expected_out):
                os.rename(expected_out, abs_pdf)
            print(f"  ✓ Converted successfully via LibreOffice.")
            return
    except Exception:
        pass

    # Windows Native Fallback
    try:
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(abs_pptx, WithWindow=False)
        presentation.SaveAs(abs_pdf, 32)
        presentation.Close()
        print(f"  ✓ Converted successfully via MS PowerPoint.")
    except Exception as e:
        print(f"  ⚠ Conversion failed: {e}")
        sys.exit(1)
    finally:
        try:
            if 'powerpoint' in locals() and powerpoint:
                powerpoint.Quit()
            pythoncom.CoUninitialize()
        except:
            pass


# ──────────────────────────────────────────────
# Part 1: python-pptx Extraction
# ──────────────────────────────────────────────
def extract_shapes(shapes, slide_info, ctx):
    """Recursively extract text, images, and tables from shapes."""
    for shape in shapes:
        if shape.shape_type == 6:
            extract_shapes(shape.shapes, slide_info, ctx)
            continue

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text and text != slide_info.get("title"):
                    slide_info["text_content"].append(text)

        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            ctx["image_counter"] += 1
            img = shape.image
            ext = img.content_type.split("/")[-1]
            
            # Auto-convert WMF/EMF Vector Graphics to PNG for Vision API compatibility
            if ext in ["x-wmf", "wmf", "x-emf", "emf"]:
                ext = "png"
                filename = f"slide_{ctx['slide_idx']:03d}_img_{ctx['image_counter']:03d}.{ext}"
                filepath = ctx["images_dir"] / filename
                
                try:
                    import io
                    from PIL import Image
                    wmf_img = Image.open(io.BytesIO(img.blob))
                    wmf_img.save(filepath, format="PNG")
                    content_type = "image/png"
                    size_bytes = os.path.getsize(filepath)
                except Exception as e:
                    print(f"  ⚠ Failed to convert WMF/EMF to PNG. Saving raw blob. Error: {e}")
                    # Fallback to original
                    filename = f"slide_{ctx['slide_idx']:03d}_img_{ctx['image_counter']:03d}.x-wmf"
                    filepath = ctx["images_dir"] / filename
                    with open(filepath, "wb") as f:
                        f.write(img.blob)
                    content_type = img.content_type
                    size_bytes = len(img.blob)
            else:
                if ext == "jpeg":
                    ext = "jpg"
                filename = f"slide_{ctx['slide_idx']:03d}_img_{ctx['image_counter']:03d}.{ext}"
                filepath = ctx["images_dir"] / filename
                with open(filepath, "wb") as f:
                    f.write(img.blob)
                content_type = img.content_type
                size_bytes = len(img.blob)

            slide_info["images"].append({
                "image_id": f"img_{ctx['slide_idx']:03d}_{ctx['image_counter']:03d}",
                "filename": filename,
                "file_path": str(filepath),
                "content_type": content_type,
                "size_bytes": size_bytes,
            })

        if shape.has_table:
            table = shape.table
            headers = [cell.text.strip() for cell in table.rows[0].cells] if table.rows else []
            rows = []
            for row_idx, row in enumerate(table.rows):
                if row_idx == 0:
                    continue 
                row_data = [cell.text.strip() for cell in row.cells]
                rows.append(row_data)

            slide_info["tables"].append({
                "headers": headers,
                "rows": rows,
                "total_rows": len(table.rows),
                "total_cols": len(table.columns),
            })

def extract_with_pptx(pptx_path: str) -> dict:
    """Extract structured content from PPTX using python-pptx."""
    prs = Presentation(pptx_path)
    images_dir = OUTPUT_DIR / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    slides_data = []
    ctx = {"image_counter": 0, "images_dir": images_dir, "slide_idx": 0}

    for slide_idx, slide in enumerate(prs.slides, 1):
        ctx["slide_idx"] = slide_idx
        slide_info = {
            "slide_number": slide_idx,
            "title": None,
            "text_content": [],
            "speaker_notes": None,
            "images": [],
            "tables": [],
        }

        if slide.shapes.title:
            slide_info["title"] = slide.shapes.title.text.strip()

        extract_shapes(slide.shapes, slide_info, ctx)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_info["speaker_notes"] = notes

        slides_data.append(slide_info)

    return {
        "metadata": {
            "source_file": Path(pptx_path).name,
            "total_slides": len(prs.slides),
            "total_images_extracted": ctx["image_counter"],
            "slides_with_notes": sum(1 for s in slides_data if s["speaker_notes"]),
        },
        "slides": slides_data,
    }


# ──────────────────────────────────────────────
# Part 2: Azure Content Understanding Analysis (PDF)
# ──────────────────────────────────────────────
def analyze_with_content_understanding(pdf_path: str) -> dict | None:
    """Send the PDF to Azure Content Understanding for AI-powered analysis."""
    if not CONTENT_UNDERSTANDING_ENDPOINT or not CONTENT_UNDERSTANDING_KEY:
        print("[INFO] Content Understanding credentials not configured. Skipping AI analysis.")
        return None

    print(f"\n[CU] Sending PDF to Content Understanding for analysis...")

    analyze_binary_url = (
        f"{CONTENT_UNDERSTANDING_ENDPOINT}/contentunderstanding"
        f"/analyzers/design_document_converter:analyzeBinary?api-version={API_VERSION}"
    )

    with open(pdf_path, "rb") as f:
        file_bytes = f.read()

    headers = {
        "Ocp-Apim-Subscription-Key": CONTENT_UNDERSTANDING_KEY,
        "Content-Type": "application/pdf", 
    }

    response = requests.post(analyze_binary_url, headers=headers, data=file_bytes)

    if response.status_code not in (200, 202):
        print(f"[CU] Error submitting file: {response.status_code}")
        print(f"[CU] Trying base64 URL method as fallback...")
        return _analyze_with_base64_fallback(pdf_path)

    operation_url = response.headers.get("Operation-Location")
    if not operation_url:
        if response.status_code == 200:
            return response.json()
        return None

    print(f"[CU] Analysis submitted. Polling for results...")
    return _poll_for_result(operation_url)


def _analyze_with_base64_fallback(pdf_path: str) -> dict | None:
    analyze_url = (
        f"{CONTENT_UNDERSTANDING_ENDPOINT}/contentunderstanding"
        f"/analyzers/design_document_converter:analyze?api-version={API_VERSION}"
    )

    with open(pdf_path, "rb") as f:
        file_bytes = f.read()

    b64_content = base64.b64encode(file_bytes).decode("utf-8")
    data_uri = f"data:application/pdf;base64,{b64_content}"

    headers = {
        "Ocp-Apim-Subscription-Key": CONTENT_UNDERSTANDING_KEY,
        "Content-Type": "application/json",
    }
    payload = {"inputs": [{"data": data_uri}]}

    response = requests.post(analyze_url, headers=headers, json=payload)
    if response.status_code not in (200, 202):
        print(f"[CU] Fallback also failed: {response.status_code}")
        return None

    operation_url = response.headers.get("Operation-Location")
    if operation_url:
        return _poll_for_result(operation_url)
    return response.json() if response.status_code == 200 else None


def _poll_for_result(operation_url: str, max_retries: int = 60, interval: int = 3) -> dict | None:
    headers = {"Ocp-Apim-Subscription-Key": CONTENT_UNDERSTANDING_KEY}
    for attempt in range(max_retries):
        time.sleep(interval)
        response = requests.get(operation_url, headers=headers)
        if response.status_code != 200:
            continue
        result = response.json()
        status = result.get("status", "")
        if status == "Succeeded":
            print(f"[CU] Analysis completed successfully!")
            return result
        elif status in ("Failed", "Cancelled"):
            print(f"[CU] Analysis {status}: {result.get('error', 'Unknown error')}")
            return None
    return None


# ──────────────────────────────────────────────
# Part 2.5: GPT-4o Vision Image Analysis
# ──────────────────────────────────────────────
# ──────────────────────────────────────────────
# Part 2.5: GPT-4o Vision Image Analysis
# ──────────────────────────────────────────────
def analyze_images_with_vision(slides_data: list) -> list:
    """Send extracted images to GPT-4o to filter out decorative images and describe architecture diagrams."""
    if not AZURE_OPENAI_ENDPOINT or not AZURE_OPENAI_KEY:
        print("[INFO] Azure OpenAI credentials not configured. Skipping Vision analysis.")
        return slides_data

    print(f"\n[VISION] Analyzing extracted images with GPT-4o...")
    
    # FIX 1: Use the new OpenAI-compatible endpoint from your Foundry snippet
    base_endpoint = AZURE_OPENAI_ENDPOINT.split("/openai")[0].rstrip("/")
    url = f"{base_endpoint}/openai/v1/chat/completions"

    # FIX 2: Added Authorization header to match OpenAI standard
    headers = {
        "api-key": AZURE_OPENAI_KEY,
        "Authorization": f"Bearer {AZURE_OPENAI_KEY}",
        "Content-Type": "application/json"
    }

    vision_prompt = (
        "You are an Expert Azure Solutions Architect. Analyze this image extracted from a pre-sales presentation.\n"
        "1. If this image is a company logo, a decorative background, a transition slide, or a generic stock photo, reply ONLY with the word: DECORATIVE.\n"
        "2. If this image is a technical architecture diagram, a data workflow, or a system component map, provide a highly detailed, professional text description of the data flow, the components involved, and how they connect.\n"
        "Do not include conversational filler."
    )

    for slide in slides_data:
        for img in slide.get("images", []):
            filepath = img.get("file_path")
            if not filepath or not os.path.exists(filepath):
                continue

            # Convert image to base64
            with open(filepath, "rb") as f:
                base64_image = base64.b64encode(f.read()).decode('utf-8')
            
            mime_type = img.get("content_type", "image/png")

            payload = {
                "model": AZURE_OPENAI_DEPLOYMENT_NAME,  # FIX 3: Add model name into the payload
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": vision_prompt},
                            {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{base64_image}"}}
                        ]
                    }
                ],
                "max_tokens": 800,
                "temperature": 0.0
            }

            try:
                response = requests.post(url, headers=headers, json=payload)
                if response.status_code == 200:
                    description = response.json()["choices"][0]["message"]["content"].strip()
                    img["ai_description"] = description
                    if description == "DECORATIVE":
                        print(f"  - {img['filename']}: Filtered as DECORATIVE")
                    else:
                        print(f"  - {img['filename']}: Architecture Described ({len(description)} chars)")
                else:
                    print(f"  ⚠ Vision analysis failed for {img['filename']}: {response.status_code}")
                    print(f"    Reason: {response.text}")
            except Exception as e:
                print(f"  ⚠ Vision error for {img['filename']}: {e}")

    return slides_data


# ──────────────────────────────────────────────
# Part 3: Merge Results
# ──────────────────────────────────────────────
def merge_extraction_results(pptx_data: dict, cu_data: dict | None) -> dict:
    result = {
        "metadata": pptx_data["metadata"],
        "extraction_sources": ["python-pptx", "azure-openai-vision"],
        "slides": pptx_data["slides"],
        "content_understanding": None,
    }

    if cu_data:
        result["extraction_sources"].append("content-understanding")
        cu_result = cu_data.get("result", cu_data)
        contents = cu_result.get("contents", [])
        if contents:
            markdown_content = contents[0].get("markdown", "")
            result["content_understanding"] = {
                "markdown": markdown_content,
                "markdown_length": len(markdown_content),
            }
    return result


# ──────────────────────────────────────────────
# Part 3: Phase 2 - Orchestrator (Table of Contents)
# ──────────────────────────────────────────────
def generate_table_of_contents(extraction_payload: dict) -> dict:
    """Uses GPT-4.1 to analyze the extracted payload and generate a structured JSON outline."""
    if not AGENT_OPENAI_ENDPOINT or not AGENT_OPENAI_KEY:
        print("[INFO] Agent credentials not configured. Skipping Phase 2.")
        return None

    print(f"\n[PHASE 2] Orchestrator is designing the document structure...")
    
    base_endpoint = AGENT_OPENAI_ENDPOINT.split("/openai")[0].rstrip("/")
    url = f"{base_endpoint}/openai/v1/chat/completions"

    headers = {
        "api-key": AGENT_OPENAI_KEY,
        "Authorization": f"Bearer {AGENT_OPENAI_KEY}",
        "Content-Type": "application/json"
    }

    # OLD PROMPT - Preserved per user request
    # orchestrator_prompt = (
    #     "You are an Expert Enterprise Solutions Architect. Your task is to analyze an extracted solution design presentation "
    #     "(provided as a combined JSON payload) and design a comprehensive Table of Contents (Outline) for a formal Microsoft Word design document.\n\n"
    #     "INSTRUCTIONS:\n"
    #     "1. Analyze the provided payload. Read the text, speaker_notes, and ai_descriptions to understand the deep technical context.\n"
    #     "2. Do NOT map slides 1-to-1 to document sections. Instead, logically group related slides into major architectural topics.\n"
    #     "3. Automatically inject necessary standard architectural sections (like Risks or Exec Summary) even if only briefly implied. DO NOT generate Section 1.0 (Document Control). Start ALL section numbering at 2.0.\n"
    #     "4. IGNORE any slides related to 'Q&A', 'Questions and Answers', 'Next Steps', 'Thank You', or 'Appendices'. DO NOT include them in the document outline.\n"
    #     "5. Output the outline strictly as a JSON object matching this schema:\n"
    #     "{\n"
    #     '  "client_name": "string (Extract ONLY the client name from the title slide. If not explicitly found, try to infer it from the context or leave blank. Do NOT include project-related words.)",\n'
    #     '  "project_title": "string (Extract ONLY the core project name or technology from the title slide. Do NOT add words like \'Solution Design Document\', \'Proposal\', or \'Document\')",\n'
    #     '  "sections": [\n'
    #     "    {\n"
    #     '      "section_number": "2.0",\n'
    #     '      "section_title": "string",\n'
    #     '      "mapped_slides": [int], \n'
    #     '      "generation_instructions": "string" \n'
    #     "    }\n"
    #     "  ]\n"
    #     "}"
    # )

    # orchestrator_prompt = (
    #     "You are an Expert Enterprise Solutions Architect from Enfrasys Solutions. Your task is to analyze an extracted solution design presentation "
    #     "(provided as a combined JSON payload) and produce a detailed Table of Contents (Outline) for a formal Microsoft Word Solution Design Document that follows the Enfrasys standard structure.\n\n"
    #     "INSTRUCTIONS:\n"
    #     "1. Analyze the payload deeply: read slide text, speaker notes, and diagram descriptions to fully understand the technical context and design decisions made.\n"
    #     "2. Do NOT map slides 1-to-1 to document sections. Group related slides into the standard Enfrasys architectural sections listed below.\n"
    #     "3. REQUIRED SECTION ORDER — Follow this exact pattern, adapting section titles to the technology in the presentation:\n"
    #     "   2.0 Executive Summary\n"
    #     "       Sub-sections: Project Overview | Document Purpose | Document Audience\n"
    #     "   3.0 Network Design and Decision\n"
    #     "       Sub-sections: Network Connectivity Overview | [Technology] Data Gateway (with VM spec table if a gateway VM exists)\n"
    #     "   4.0 Roles in [Technology]\n"
    #     "       Sub-sections: [Technology] Administrator Role (include: H3 for Assigning the Role + H3 for Admin Portal capabilities table) | [Technology]-native Roles (H3 for Workspace-level Roles with capabilities table + H3 for Item-level Roles with permissions table) | [Client] Workspace Role and Access (client-specific role mapping table)\n"
    #     "   5.0 [Technology] Design and Decision  ← This is the LARGEST section. It MUST cover ALL of the following sub-sections:\n"
    #     "       5.x [Topic] Workflow Design Considerations (table of design considerations with IDs and descriptions)\n"
    #     "       5.x [Topic] Workflow Design Decisions\n"
    #     "       5.x Access Design Considerations (H3 per access method: e.g. Conditional Access, Private Link, MFA)\n"
    #     "       5.x Access Design Decisions (table listing final access policy decisions)\n"
    #     "       5.x Identity and Access Management\n"
    #     "       5.x Virtual Machine (VM) Size (table with OS, .NET, vCPU, RAM specs for any gateway/VM component)\n"
    #     "       5.x Resource Organization Design (H3: Management Groups | H3: Naming and Tagging → H4 tables: Azure Naming Convention, [Tech] Naming Convention, Application Acronyms, Suggested Tags)\n"
    #     "       5.x Governance Consideration (H3: Governance Disciplines with Cost Management Tools table | H3: Violation Triggers and Actions table)\n"
    #     "       5.x Multi-Factor Authentication (if mentioned in the slides)\n"
    #     "   6.0 Security Design and Decision\n"
    #     "       Sub-sections: Azure NSG Overview (H3: Inbound Rules table if applicable | H3: Outbound Rules table) | Encryption Design if applicable\n"
    #     "   7.0 [Unique Technology Section — e.g. AutoML, Machine Learning, Migration]\n"
    #     "       Sub-sections: Technical Overview | Approach/Paths (H3 per path) | Process Flow (H3 per phase) | Design Decisions (table with IDs)\n"
    #     "   8.0 Appendix\n"
    #     "       Sub-sections: Appendix 1 Computing | Appendix 2 Network | Appendix 3 Identity & Security | Appendix 4 Logging & Monitoring | Appendix 5 Cloud Governance | Appendix 6 [Tech-specific e.g. OPDG] | Appendix 7 [Tech-specific e.g. Machine Learning]\n"
    #     "       EACH appendix sub-section MUST have three H3 sub-sub-sections: 'Introduction/Prerequisites' (reference documentation links), 'Limits and Boundaries' (quota/limit links), 'Others' (additional references).\n"
    #     "4. CONTENT PLACEMENT RULES:\n"
    #     "   - VM specifications belong in the Platform Design section (5.x Virtual Machine Size), NOT in the Network section.\n"
    #     "   - NSG rules (Inbound/Outbound) belong in the Security Design section ONLY. If referenced elsewhere, refer to the section number rather than repeating the table.\n"
    #     "   - Conditional Access and MFA decisions belong in the Platform Design section (Access Design). Security section covers NSG/encryption only.\n"
    #     "   - Naming conventions belong in Resource Organization Design ONLY. Do not repeat them in the Appendix.\n"
    #     "5. Do NOT create standalone top-level sections for 'Microsoft Best Practices', 'Governance', 'Limitations', or 'Naming Conventions' — these are integrated sub-sections only.\n"
    #     "6. Do NOT generate Section 1.0 (Document Sign Off). Start ALL section numbering at 2.0.\n"
    #     "7. IGNORE slides for 'Q&A', 'Questions and Answers', 'Next Steps', 'Thank You', or 'Appendices'.\n"
    #     "8. TARGET SCOPE: Produce 8-10 top-level sections as defined above, but ensure DEEP granularity in the sub-sections to produce a comprehensive, lengthy document.\n"
    #     "9. Output the outline strictly as a JSON object:\n"
    #     "{\n"
    #     '  "client_name": "string (Extract exactly as written on the title slide. Do not expand abbreviations.)",\n'
    #     '  "client_name_full": "string | null (Full name only if 100% certain. Otherwise null.)",\n'
    #     '  "project_title": "string (Core technology or project name only. Do not add words like Solution Design Document.)",\n'
    #     '  "sections": [\n'
    #     "    {\n"
    #     '      "section_number": "2.0",\n'
    #     '      "section_title": "string",\n'
    #     '      "mapped_slides": [int],\n'
    #     '      "generation_instructions": "string (Be explicit: list every H2 and H3 sub-section the Writer must produce, every table required (with column names), every design decision to document, and which image to embed. State clearly which content belongs only in this section and must NOT be repeated elsewhere.)"\n'
    #     "    }\n"
    #     "  ]\n"
    #     "}"
    # )

    orchestrator_prompt = (
        "You are an Expert Enterprise Solutions Architect from Enfrasys Solutions. Your task is to analyze an extracted solution design presentation "
        "(provided as a combined JSON payload) and produce a detailed Table of Contents (Outline) for a formal Microsoft Word Solution Design Document that follows the Enfrasys standard structure.\n\n"
        "INSTRUCTIONS:\n"
        "1. Analyze the payload deeply: read slide text, speaker notes, and diagram descriptions to fully understand the technical context and design decisions made.\n"
        "2. Group related slides into the standard Enfrasys architectural sections listed below. Do NOT map slides 1-to-1.\n"
        "3. REQUIRED SECTION ORDER — Follow this exact pattern, including strict decimal numbering (e.g., 2.1, 2.2) for all sub-sections. Adapt section titles to the technology in the presentation:\n"
        "   2.0 Executive Summary\n"
        "       Sub-sections: 2.1 Project Overview | 2.2 Document Purpose | 2.3 Document Audience\n"
        "   3.0 [Technology] Overview\n"
        "       Sub-sections: 3.1 [Core Technology 1] | 3.2 [Core Technology 2] (Force the Writer to generate boilerplate definitions for the core Microsoft technologies used, e.g., Fabric, M365, Entra ID).\n"
        "   4.0 Network Design and Decision\n"
        "       Sub-sections: 4.1 Network Connectivity Overview | 4.2 [Technology] Data Gateway (with VM spec table if a gateway VM exists)\n"
        "   5.0 Roles in [Technology]\n"
        "       Sub-sections: 5.1 [Technology] Administrator Role | 5.2 [Technology]-native Roles | 5.3 [Client] Workspace Role and Access\n"
        "   6.0 [Technology] Design and Decision  ← This is the LARGEST section. It MUST cover ALL of the following sub-sections:\n"
        "       6.1 [Topic] Workflow Design Considerations (table of design considerations with IDs and descriptions)\n"
        "       6.2 [Topic] Workflow Design Decisions\n"
        "       6.3 Access Design Considerations and Decisions (H3 per access method, plus decision table)\n"
        "       6.4 Virtual Machine (VM) Size (table with OS, .NET, vCPU, RAM specs for any gateway/VM component)\n"
        "       6.5 Resource Organization Design (H3: Management Groups | H3: Naming and Tagging → tables: Azure Naming, [Tech] Naming, Acronyms)\n"
        "       6.6 Governance Consideration (H3: Governance Disciplines with Cost Management table | H3: Violation Triggers and Actions table)\n"
        "   7.0 Security Design and Decision\n"
        "       Sub-sections: 7.1 Azure NSG Overview (Inbound/Outbound Rules tables) | 7.2 Encryption Design\n"
        "   8.0 Deployment & Migration Approach [If Applicable]\n"
        "       Sub-sections: 8.1 Pre-Migration/Setup | 8.2 Pilot | 8.3 Production | 8.4 Post-Migration/Support. (Instruct the Writer to generate 'Action By' task tables and D-7/D0/D+1 timelines here).\n"
        "   9.0 Appendix\n"
        "       Sub-sections: 9.1 Appendix 1 Computing | 9.2 Appendix 2 Network | 9.3 Appendix 3 Identity & Security | 9.4 Appendix 4 Logging & Monitoring | 9.5 Appendix 5 Cloud Governance | 9.6 Appendix 6 [Tech-specific]\n"
        "       EACH appendix sub-section MUST have: 'Introduction/Prerequisites', 'Limits and Boundaries', 'Others'.\n"
        "4. CONTENT PLACEMENT RULES:\n"
        "   - NSG rules belong in Security Design ONLY.\n"
        "   - VM specifications belong in the Design and Decision section ONLY.\n"
        "   - Naming conventions belong in Resource Organization Design ONLY.\n"
        "5. Do NOT create standalone top-level sections for 'Microsoft Best Practices', 'Governance', or 'Naming Conventions'.\n"
        "6. Do NOT generate Section 1.0 (Document Sign Off). Start ALL section numbering at 2.0.\n"
        "7. Output the outline strictly as a JSON object matching this schema:\n"
        "{\n"
        '  "client_name": "string (Exact client name from title slide)",\n'
        '  "client_name_full": "string | null",\n'
        '  "project_title": "string",\n'
        '  "sections": [\n'
        "    {\n"
        '      "section_number": "2.0",\n'
        '      "section_title": "string",\n'
        '      "mapped_slides": [int],\n'
        '      "generation_instructions": "string (List every H2/H3 sub-section with its exact decimal number (e.g., 2.1, 6.3), required tables WITH column names, explicit instructions to add boilerplate definitions, and instructions to assign Enfrasys vs. Client responsibilities.)"\n'
        "    }\n"
        "  ]\n"
        "}"
    )

    payload = {
        "model": ORCHESTRATOR_DEPLOYMENT,
        "response_format": { "type": "json_object" },
        "messages": [
            {"role": "system", "content": orchestrator_prompt},
            {"role": "user", "content": json.dumps(extraction_payload, default=str)}
        ],
        "temperature": 0.2
    }

    try:
        response = requests.post(url, headers=headers, json=payload)
        if response.status_code == 200:
            content = response.json()["choices"][0]["message"]["content"]
            return json.loads(content)
        else:
            print(f"  ⚠ Orchestrator failed: {response.status_code} - {response.text}")
            return None
    except Exception as e:
        print(f"  ⚠ Orchestrator error: {e}")
        return None

# ──────────────────────────────────────────────
# Part 4: Phase 3 - Writer (Drafting the Document)
# ──────────────────────────────────────────────
def write_document_sections(toc: dict, extraction_payload: dict) -> str:
    """Loops through the TOC, grabs the relevant slide data, and writes the final document."""
    print(f"\n[PHASE 3] Writer Agent is drafting the final document sections...")
    
    base_endpoint = AGENT_OPENAI_ENDPOINT.split("/openai")[0].rstrip("/")
    url = f"{base_endpoint}/openai/v1/chat/completions"
    headers = {
        "api-key": AGENT_OPENAI_KEY,
        "Authorization": f"Bearer {AGENT_OPENAI_KEY}",
        "Content-Type": "application/json"
    }

    final_document_markdown = f"# {toc.get('client_name', '')} {toc.get('project_title', 'Solution Design Document')}".strip() + "\n\n"

    # OLD PROMPT - Preserved
    # writer_system_prompt = """
    #     You are an Expert Technical Writer and Enterprise Architect. 
    #     Your job is to write a highly professional, formal section of a Solution Design Document based on specific slide data. 
    #     Maintain a formal, authoritative, and extremely concise tone. 
    #     CRITICAL: Do NOT wrap your entire response in a ```markdown code block. Just return the raw text.
    # ...
    # """

    # writer_system_prompt = """
    #     You are a Lead Enterprise Architect from Enfrasys Solutions writing a formal Solution Design Document. Your output will be converted directly to a Microsoft Word document for a client.
    #     Write in authoritative, professional language ("We recommend", "[CLIENT_NAME] shall implement"). Be specific — reference the actual client name, technology components, and design decisions from the slide data.
    #     The client expects a comprehensive, highly detailed architectural document. Expand fully on the technical reasoning, configuration specifics, and business impact for every decision.
    #     CRITICAL: Do NOT wrap your entire response in a ```markdown code block. Just return the raw text.

    #     -- MARKDOWN HEADING RULES --
    #     - Main Sections MUST use `#` (e.g., `# 5.0 Fabric Design and Decision`).
    #     - Sub-sections MUST use `##` (e.g., `## 5.1 Data Workflow Design Considerations`).
    #     - Sub-sub-sections MUST use `###` (e.g., `### 5.7.1 Azure Management Group`).
    #     - The H1 main section opening paragraph must be 1-2 sentences only. Detail goes in H2/H3 sub-sections.

    #     -- H3 DEPTH RULES (IMPORTANT) --
    #     Use ### (H3) inside H2 sections for distinct named sub-components. Required patterns:
    #     - Administrator Role H2 → ### Assigning [Role Name] + ### The [Technology] Admin Portal (with capabilities table)
    #     - Fabric-native Roles H2 → ### Workspace-level Roles (capabilities table) + ### Item-level Roles (permissions table)
    #     - Access Design H2 → ### per access method (e.g. ### Conditional Access / ### Private Link Consideration)
    #     - Resource Organization H2 → ### Azure Management Group + ### Naming and Tagging (→ H4 tables: naming convention, tech-specific naming, acronyms, suggested tags)
    #     - Governance H2 → ### Azure Governance Discipline (→ H4: Cost Management Tools table) + ### Violation Triggers and Actions (table)
    #     - Workflow/ML H2 with multiple paths → ### per path (e.g. ### AutoML Path, ### Manual ML Path, ### Final Output)
    #     - Appendix sub-sections → ### Introduction/Prerequisites + ### Limits and Boundaries + ### Others (each with bullet-list reference links)

    #     -- STRICT IMAGE RULES --
    #     - Select only architecture diagrams, data flow diagrams, network topologies, or access/security diagrams. Never use logos or decorative slides.
    #     - Each UNIQUE image file_path may only be embedded ONCE across the entire document. Do NOT embed the same file_path in two different sections.
    #     - You MUST embed all relevant architecture diagrams provided in the JSON for the current section to provide maximum visual context.
    #     - Syntax: `![](file_path)` — use the exact path from the JSON. Do NOT invent paths.

    #     -- TABLE RULES --
    #     Use Markdown tables as the primary data carrier. Required tables (produce these whenever the slide data contains this information):
    #     - Administrator capabilities: columns = Capability | Description
    #     - Workspace-level role capabilities: columns = Permission | Admin | Contributor | Member | Viewer
    #     - Client workspace role mapping: columns = No. | Workspace Role | Suggested Role | [Client] Personnel
    #     - Design Considerations: columns = ID | Description | Workload Type
    #     - Design Decisions: columns = No. | Design Decision | Decision
    #     - Access/Conditional Access decisions: columns = No. | Policy | Decision
    #     - VM specifications: columns = Component | Specification
    #     - NSG rules: columns = Name | Priority | Source | Source Ports | Destination | Destination Ports | Protocol | Access
    #     - Azure Naming Convention: columns = Resource | Abbreviation | Example
    #     - Technology-specific Naming: columns = No. | Item | Naming Convention Example
    #     - Application Acronyms: columns = Acronym | Full Name
    #     - Suggested Tags: columns = Tag Key | Description | Example Value
    #     - Cost Management Tools: columns = Azure Tool | Description | Cost Management Discipline
    #     - Violation Triggers and Actions: columns = Trigger | Action
    #     The paragraph before each table must be 1-2 sentences only explaining the table's purpose.

    #     -- NO-REPETITION RULES (CRITICAL) --
    #     1. Each table must appear EXACTLY ONCE across the full document. If a table was in a main section, the Appendix must NOT repeat it — instead, write "Refer to Section X.X for [table name]."
    #     2. NSG rules belong ONLY in the Security Design section. Do not put NSG tables in the Network section or Appendix.
    #     3. VM specifications belong ONLY in the Platform Design section (Virtual Machine Size sub-section). Do not repeat in Network or Appendix.
    #     4. Naming convention tables belong ONLY in the Resource Organization Design sub-section. Reference via section number in the Appendix.
    #     5. Access/Conditional Access decisions belong ONLY in the Platform Design section. Do not repeat in Security section.

    #     -- WRITING STYLE RULES --
    #     1. DECISION-FOCUSED: State the design decision first, then explain WHY with client-specific context.
    #        Example: "[CLIENT_NAME] shall implement Conditional Access for all Fabric users. Private Link was evaluated but deferred due to licensing constraints."
    #     2. INTEGRATE BEST PRACTICES: Weave best practices into the prose as justification — never create a 'Microsoft Best Practices' heading.
    #     3. CONTENT DEPTH: Each H2 sub-section MUST contain 3-6 highly detailed paragraphs covering the technical rationale, configuration steps, and operational impact. Each H3 needs 2-3 paragraphs plus its table if applicable. Do not be overly brief.
    #     4. APPENDIX PATTERN: Each appendix entry (e.g., Appendix 1 Computing) must have exactly three H3 sub-sections:
    #        ### Introduction/Prerequisites — bullet list of official Microsoft documentation links
    #        ### Limits and Boundaries — bullet list of Azure quota/limit reference links
    #        ### Others — bullet list of supplementary reference links
    #     5. NO GENERIC FILLER: Every sentence must be grounded in the client's specific context or a concrete design decision from the slides.
    # """

    writer_system_prompt = """
        You are a Lead Enterprise Architect from Enfrasys Solutions writing a formal Solution Design Document. Your output will be converted directly to a Microsoft Word document for a client.
        Write in authoritative, professional language ("Enfrasys recommends", "[CLIENT_NAME] shall implement"). Be specific — reference the actual client name, technology components, and design decisions from the slide data.
        
        CRITICAL: Do NOT wrap your entire response in a ```markdown code block. Just return the raw text.

        -- CONSULTING EXPANSION RULES (CRITICAL) --
        1. STATIC BOILERPLATE: Whenever you introduce a major Microsoft technology (e.g., Microsoft Fabric, Exchange Online, Entra ID, Power BI), you MUST provide a standard, formal definition of that technology before discussing the client's specific design. Do not assume the reader knows what the technology does.
        2. RESPONSIBILITY ASSIGNMENT: In any section discussing tasks, migrations, or rollouts, you MUST explicitly assign ownership. Use "Enfrasys Solutions" for vendor tasks and "[CLIENT_NAME]" for client tasks.
        3. TIMELINES (If Applicable): If the slides mention migration waves or rollouts, automatically structure them into formal consulting timelines (e.g., Pre-Migration (D-7), Cutover (D0), Post-Migration Support (D+1)).

        -- MARKDOWN HEADING RULES --
        - Main Sections MUST use `#` (e.g., `# 5.0 Fabric Design and Decision`).
        - Sub-sections MUST use `##` (e.g., `## 5.1 Data Workflow Design Considerations`).
        - Sub-sub-sections MUST use `###` (e.g., `### 5.7.1 Azure Management Group`).
        - The H1 main section opening paragraph must be 1-2 sentences only. Detail goes in H2/H3 sub-sections.

        -- TABLE RULES (STRICT SCHEMA) --
        Use Markdown tables to present all technical data. You must use the exact columns below when generating these tables:
        - Design Considerations: columns = ID | Description | Workload Type
        - Design Decisions: columns = No. | Design Decision | Decision
        - Task/Migration Rollout: columns = Item | Activities | Action By | Status (Assign 'Action By' to Enfrasys or [CLIENT_NAME])
        - Administrator capabilities: columns = Capability | Description
        - Workspace-level role capabilities: columns = Permission | Admin | Contributor | Member | Viewer
        - Client workspace role mapping: columns = No. | Workspace Role | Suggested Role | [CLIENT_NAME] Personnel
        - Access decisions: columns = No. | Policy | Decision
        - VM specifications: columns = Component | Specification
        - NSG rules: columns = Name | Priority | Source | Source Ports | Destination | Destination Ports | Protocol | Access
        - Azure Naming Convention: columns = Resource | Abbreviation | Example
        - Cost Management Tools: columns = Azure Tool | Description | Cost Management Discipline

        -- STRICT IMAGE RULES --
        - Select only architecture diagrams, data flow diagrams, network topologies, or access/security diagrams.
        - Each UNIQUE image file_path may only be embedded ONCE across the entire document.
        - Syntax: `![](file_path)` — use the exact path from the JSON. Do NOT invent paths. Wrap text around the image explaining what it shows (e.g., "Figure 1 below illustrates the data gateway architecture...").

        -- NO-REPETITION RULES --
        1. Each table must appear EXACTLY ONCE across the full document. If a table was in a main section, the Appendix must NOT repeat it.
        2. NSG rules belong ONLY in the Security Design section. 
        3. VM specifications belong ONLY in the Platform Design section.
        4. Naming convention tables belong ONLY in the Resource Organization Design sub-section.

        -- APPENDIX FORMATTING --
        Each appendix entry MUST have exactly these three H3 sub-sections:
        ### Introduction/Prerequisites (Provide real Microsoft Docs links)
        ### Limits and Boundaries (Provide real Azure quota links)
        ### Others (Supplementary links)
    """

    # Check if we should use the Foundry Assistant or standard REST API
    use_assistant = bool(AGENT_ASSISTANT_ID and 'AzureOpenAI' in globals())
    
    if use_assistant:
        print(f"  -> Connecting to Microsoft Foundry Agent (ID: {AGENT_ASSISTANT_ID})...")
        try:
            client = AzureOpenAI(
                api_key=AGENT_OPENAI_KEY,  
                api_version="2024-05-01-preview", # Assistant API version
                azure_endpoint=AGENT_OPENAI_ENDPOINT
            )
            # 1. Dynamically sync the Foundry Prompt to match our Python code!
            client.beta.assistants.update(
                assistant_id=AGENT_ASSISTANT_ID,
                instructions=writer_system_prompt
            )
            print("  ✓ Foundry Agent Instructions successfully synced with Source Code.")
        except Exception as e:
            print(f"  ⚠ Failed to initialize or update Assistant: {e}")
            use_assistant = False

    for section in toc.get("sections", []):
        sec_num = section.get("section_number")
        sec_title = section.get("section_title")
        instructions = section.get("generation_instructions")
        mapped_slides = section.get("mapped_slides", [])

        print(f"  -> Drafting {sec_num} {sec_title}...")

        # Gather only the data for the mapped slides
        relevant_slide_data = [s for s in extraction_payload.get("slides", []) if s.get("slide_number") in mapped_slides]

        client_name = toc.get("client_name", "")
        client_name_full = toc.get("client_name_full", "")
        
        context_str = f"Client Acronym: {client_name}\n"
        if client_name_full:
            context_str += f"Client Full Name: {client_name_full} (Use this full name in the content when appropriate)\n"

        user_prompt = (
            f"Write section '{sec_num} {sec_title}' for the design document.\n"
            f"{context_str}\n"
            f"Instructions from Orchestrator: {instructions}\n\n"
            f"Here is the raw slide data (text, speaker notes, and diagram descriptions):\n"
            f"{json.dumps(relevant_slide_data, default=str)}"
        )

        drafted_text = None

        if use_assistant:
            # === NEW: FIRE UP THE FOUNDRY AGENT THREAD ===
            try:
                # 2. Create a clean Thread
                thread = client.beta.threads.create()
                
                # 3. Add the Payload Message
                client.beta.threads.messages.create(
                    thread_id=thread.id,
                    role="user",
                    content=user_prompt
                )
                
                # 4. Trigger the Run and let Web Search act natively
                run = client.beta.threads.runs.create_and_poll(
                    thread_id=thread.id,
                    assistant_id=AGENT_ASSISTANT_ID
                )
                
                if run.status == 'completed':
                    messages = client.beta.threads.messages.list(thread_id=thread.id)
                    drafted_text = messages.data[0].content[0].text.value
                else:
                    print(f"  ⚠ Agent run failed. Status: {run.status}")
                
                # 5. Destroy the Thread to prevent memory hallucination!
                client.beta.threads.delete(thread.id)
                
            except Exception as e:
                print(f"  ⚠ Foundry Agent error on {sec_num}: {e}")
        
        # === FALLBACK: STANDARD REST API ===
        if not use_assistant or not drafted_text:
            payload = {
                "model": WRITER_DEPLOYMENT,
                "messages": [
                    {"role": "system", "content": writer_system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                "temperature": 0.4
            }

            try:
                response = requests.post(url, headers=headers, json=payload)
                if response.status_code == 200:
                    drafted_text = response.json()["choices"][0]["message"]["content"].strip()
                else:
                    print(f"  ⚠ REST Writer failed for section {sec_num}: {response.text}")
            except Exception as e:
                print(f"  ⚠ REST Writer error for section {sec_num}: {e}")

        # Post-process generated text
        if drafted_text:
            if drafted_text.startswith("```"):
                drafted_text = drafted_text.split("\n", 1)[-1]
            if drafted_text.endswith("```"):
                drafted_text = drafted_text.rsplit("\n", 1)[0]
            
            drafted_text = drafted_text.strip()
            final_document_markdown += f"{drafted_text}\n\n"
        else:
            final_document_markdown += f"## {sec_num} {sec_title}\n\n*(Error generating section)*\n\n"

    return final_document_markdown


# ──────────────────────────────────────────────
# Part 5: Phase 4 - Convert to Word Document
# ──────────────────────────────────────────────
def convert_md_to_docx(md_path, docx_path, document_title, project_title, client_name, client_logo_path=None):
    """Converts the final Markdown file into a formatted Microsoft Word document and fixes tables/page breaks/footers."""
    print(f"\n[PHASE 4] Converting Markdown to Microsoft Word...")
    try:
        import pypandoc
        from docx import Document
        from docx.oxml.ns import qn        
        from docx.oxml import OxmlElement  
        
        # Point to the blank Word template
        template_file = str(SCRIPT_DIR / "enfrasys_template.docx")
        
        if os.path.exists(template_file):
            print(f"  -> Applying corporate styles from {template_file}")
            extra_args = [f'--reference-doc={template_file}']
        else:
            print("  -> No template found. Using default Word styles.")
            extra_args = []

        # Convert using Pandoc
        pypandoc.convert_file(
            str(md_path), 
            'docx', 
            outputfile=str(docx_path),
            extra_args=extra_args
        )
        print(f"  ✓ Word document generated.")

        # --- POST-PROCESSING WITH PYTHON-DOCX ---
        doc = Document(str(docx_path))
        
        # 0. TITLE PAGE GENERATION
        print(f"  -> Generating Template Title Page...")
        
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        # Extract the very first paragraph which contains the title generated by markdown (# Title)
        # We will clear it because we are building a custom title page.
        if len(doc.paragraphs) > 0:
            first_p = doc.paragraphs[0]
            first_p.clear()
        else:
            first_p = doc.add_paragraph()
            
        # Remove any page breaks before the VERY first paragraph just in case
        first_p.paragraph_format.page_break_before = False
            
        # We will insert paragraphs before the old title paragraph to build the title page.
        
        # 1. Black Bar (Simplest way: huge font with black highlight/background, or dark border)
        # Since shading XML is complex, we will format the text " " with a black background over the whole line.
        # Actually, using a paragraph with a thick bottom border is robust.
        def set_p_bottom_border(p):
            pPr = p._element.get_or_add_pPr()
            pBdr = OxmlElement('w:pBdr')
            pPr.append(pBdr)
            bottom = OxmlElement('w:bottom')
            bottom.set(qn('w:val'), 'single')
            bottom.set(qn('w:sz'), '144') # 144 = 18 pt thickness
            bottom.set(qn('w:space'), '1')
            bottom.set(qn('w:color'), '000000')
            pBdr.append(bottom)

        p_bar = first_p.insert_paragraph_before("")
        set_p_bottom_border(p_bar)
        
        # 2. Project Title Only
        p_title = first_p.insert_paragraph_before()
        p_title.paragraph_format.space_before = Pt(6)
        run = p_title.add_run(project_title)
        run.font.name = 'Segoe UI' # Standard modern font
        run.font.size = Pt(22)
        run.font.color.rgb = RGBColor(0, 0, 0)
        
        # 3. "Planning & Design Document" Subtitle
        p_sub = first_p.insert_paragraph_before()
        p_sub.paragraph_format.space_before = Pt(24)
        run = p_sub.add_run("Planning & Design Document")
        run.font.name = 'Segoe UI'
        run.font.size = Pt(22)
        run.font.color.rgb = RGBColor(0, 0, 0)
        
        # 3.5 Client Logo
        from docx.shared import Inches
        if client_logo_path and os.path.exists(client_logo_path):
            p_logo = first_p.insert_paragraph_before()
            p_logo.paragraph_format.space_before = Pt(40)
            run = p_logo.add_run()
            try:
                run.add_picture(str(client_logo_path), height=Inches(1.5))
            except Exception as e:
                print(f"  ⚠ Failed to insert client logo: {e}")
        
        # 4. Spacing before next sections
        # We need to drop down several lines for the footer-like abstract
        for _ in range(4): # Reduced to 6 to make room for logo
            first_p.insert_paragraph_before("")
            
        # 5. Abstract Divider
        p_abs_div = first_p.insert_paragraph_before("")
        def set_p_thin_bottom_border(p):
            pPr = p._element.get_or_add_pPr()
            pBdr = OxmlElement('w:pBdr')
            pPr.append(pBdr)
            bottom = OxmlElement('w:bottom')
            bottom.set(qn('w:val'), 'single')
            bottom.set(qn('w:sz'), '8') # 1 pt thickness
            bottom.set(qn('w:space'), '1')
            bottom.set(qn('w:color'), '000000')
            pBdr.append(bottom)
        set_p_thin_bottom_border(p_abs_div)
        
        # Abstract Title
        p_abs_title = first_p.insert_paragraph_before()
        p_abs_title.paragraph_format.space_before = Pt(6)
        run = p_abs_title.add_run("Abstract")
        run.font.name = 'Segoe UI'
        run.bold = True
        run.font.size = Pt(11)
        
        # Abstract Text
        p_abs_text = first_p.insert_paragraph_before()
        p_abs_text.paragraph_format.space_before = Pt(6)
        run = p_abs_text.add_run(f"This document defines the Planning and Design for {document_title} project.")
        run.font.name = 'Segoe UI'
        run.font.size = Pt(11)
        
        # Prepared By
        for _ in range(2):
            first_p.insert_paragraph_before("")
            
        p_prep = first_p.insert_paragraph_before()
        run = p_prep.add_run("Prepared By")
        run.font.name = 'Segoe UI'
        run.font.size = Pt(11)
        
        # 6. Enfrasys Logo
        enfrasys_logo = SCRIPT_DIR / "enfrasys-logo.jpg"
        if os.path.exists(enfrasys_logo):
            p_enfrasys = first_p.insert_paragraph_before()
            p_enfrasys.paragraph_format.space_before = Pt(10)
            run = p_enfrasys.add_run()
            try:
                run.add_picture(str(enfrasys_logo), width=Inches(1.5))
            except Exception as e:
                print(f"  ⚠ Failed to insert Enfrasys logo: {e}")
        
        # Add a Page Break right after the Prepared By section to move everything else to Page 2
        from docx.enum.text import WD_BREAK
        p_break_title = first_p.insert_paragraph_before("")
        p_break_title.add_run().add_break(WD_BREAK.PAGE)
        
        # ---------------------------------------------------------
        # 0.5 PAGE 2 BUILDING (Document Control & Header)
        # ---------------------------------------------------------
        print(f"  -> Generating Document Control Page & Headers...")
        
        # A) Setup Header Logos
        doc.sections[0].different_first_page_header_footer = True
        header = doc.sections[0].header
        
        # The header already has a default paragraph. Let's create a 1-row, 2-column invisible table for the logos
        htable = header.add_table(1, 2, Inches(6.0))
        # Clear any default empty paragraphs in the header so it doesn't take up extra space
        for p in header.paragraphs:
            if not p._element.getparent() == htable._element and p.text.strip() == "":
                p_elem_h = p._element
                p_elem_h.getparent().remove(p_elem_h)
                
        # Make the header table invisible
        tblPr_h = htable._tbl.tblPr
        tblBorders = OxmlElement('w:tblBorders')
        for border_name in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
            border = OxmlElement(f'w:{border_name}')
            border.set(qn('w:val'), 'none')
            tblBorders.append(border)
        tblPr_h.append(tblBorders)
        
        cell_left = htable.cell(0, 0)
        cell_right = htable.cell(0, 1)
        
        # Add Client Logo to Header Left
        if client_logo_path and os.path.exists(client_logo_path):
            p_left = cell_left.paragraphs[0]
            p_left.alignment = WD_ALIGN_PARAGRAPH.LEFT
            try:
                p_left.add_run().add_picture(str(client_logo_path), height=Inches(0.6))
            except Exception:
                pass
                
        # Add Enfrasys Logo to Header Right
        from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
        cell_right.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        if os.path.exists(enfrasys_logo):
            p_right = cell_right.paragraphs[0]
            p_right.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            try:
                p_right.add_run().add_picture(str(enfrasys_logo), width=Inches(1.5))
            except Exception:
                pass
                
        # B) Setup Page 2 Content
        # We use a brand new paragraph inserted AFTER the title page break as our anchor
        anchor_p = first_p.insert_paragraph_before("")
        p_xml = anchor_p._element
        
        def add_p_before(text="", bold=False, size=None, space_before=None):
            p = anchor_p.insert_paragraph_before()
            if space_before is not None: 
                p.paragraph_format.space_before = Pt(space_before)
            if text:
                r = p.add_run(text)
                if bold: r.bold = True
                if size: r.font.size = Pt(size)
                r.font.name = 'Segoe UI'
            return p

        # 1. Change Record
        add_p_before("Change Record", bold=True, size=14)
        
        table1 = doc.add_table(rows=3, cols=4)
        try:
            table1.style = 'Table Grid'
        except KeyError:
            pass
            
        hdr_cells = table1.rows[0].cells
        hdr_cells[0].text = 'Date'
        hdr_cells[1].text = 'Author'
        hdr_cells[2].text = 'Version'
        hdr_cells[3].text = 'Change Reference'
        
        row_cells = table1.rows[1].cells
        row_cells[0].text = '[change-record-date]'
        row_cells[1].text = '[change-record-author]'
        row_cells[2].text = '[change-record-version]'
        row_cells[3].text = '[change-record-reference]'
        
        # Move table before anchor
        p_xml.addprevious(table1._element)
        
        add_p_before("", space_before=12)
        add_p_before("Distribution List", bold=True, size=14)
        
        table2 = doc.add_table(rows=3, cols=2)
        try:
            table2.style = 'Table Grid'
        except KeyError:
            pass
            
        hdr_cells2 = table2.rows[0].cells
        hdr_cells2[0].text = 'Name'
        hdr_cells2[1].text = 'Position/ Team'
        
        row_cells2 = table2.rows[1].cells
        row_cells2[0].text = '[Client/Entity/Company Name]'
        row_cells2[1].text = '[Client/Entity/Company Position]'
        
        row_cells3 = table2.rows[2].cells
        row_cells3[0].text = ''
        row_cells3[1].text = ''
        
        p_xml.addprevious(table2._element)
        
        add_p_before("", space_before=24)
        
        def add_legal(title, text):
            p = add_p_before(title, bold=True, size=7, space_before=12)
            p.paragraph_format.left_indent = Inches(2.3)
            
            ptext = add_p_before(text, size=7)
            ptext.paragraph_format.left_indent = Inches(2.3)
            ptext.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            
        add_legal("Terms of Use", "This work may be used \"as-is\" by any interested party. You may copy, adapt, and redistribute this document for non-commercial use or for your own internal use in a commercial setting. However, you may not republish this document, nor may you publish or distribute any adaptation of this document for other than non-commercial use or your own internal use, without first obtaining express written approval from Enfrasys Solutions Sdn Bhd.")
        add_legal("Disclaimer", "The Author and Enfrasys Solutions Sdn Bhd shall have neither liability nor responsibility to any person or entity with respect to the loss or damages arising from the information contained in this work. This work may include inaccuracies or typographical errors and solely represent the opinions of the Author. Changes are periodically made to this document without notice.Due to the rapid growth of various technologies, the Author and Enfrasys Solutions Sdn Bhd cannot guarantee the accuracy of the information presented after the date of publication.")
        add_legal("Trademarks", "The names of actual companies, service marks, trademarks or products mentioned herein may be the trademarks of their respective owners. Use of terms within this work should not be regarded as affecting the validity of any trademark or service mark. Enfrasys Consulting may have patents, patent applications, trademarks, copyrights, or other intellectual property rights covering subject matter in this document. Except as expressly provided in any written license agreement from Enfrasys Solutions Sdn Bhd, the furnishing of this document does not give you any license to those items.")
        
        # Add page break to jump to Document Sign Off
        p_break2 = add_p_before("")
        p_break2.add_run().add_break(WD_BREAK.PAGE)
        
        # ---------------------------------------------------------
        # 0.6 PAGE 3 & 4 (1.0 Design Document Sign Off)
        # ---------------------------------------------------------
        p_signoff = add_p_before("1.0 Design Document Sign Off", size=15)
        try:
            p_signoff.style = 'Heading 1'
        except Exception:
            pass
            
        signoff_text_1 = f"We hereby acknowledge that the design document for {client_name} {project_title} has been reviewed, and all key aspects have been addressed satisfactorily."
        add_p_before(signoff_text_1, size=11, space_before=12)
        
        signoff_text_2 = "This document has been prepared, reviewed, accepted, and signed off by the following individuals:"
        add_p_before(signoff_text_2, size=11, space_before=12)
        add_p_before("", space_before=4) # Increased gap between Text and Table!
        
        # Hardcode black borders in case the corporate template hides 'Table Grid'
        def set_table_borders(table):
            tbl_borders = OxmlElement('w:tblBorders')
            for border_name in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
                border = OxmlElement(f'w:{border_name}')
                border.set(qn('w:val'), 'single')
                border.set(qn('w:sz'), '8')  # 1 pt border
                border.set(qn('w:color'), '000000') # Dark Black
                tbl_borders.append(border)
            table._tbl.tblPr.append(tbl_borders)
            
        # TABLE 1 (Prepared By / Verified By)
        table_sig_1 = doc.add_table(rows=1, cols=2)
        try: table_sig_1.style = 'Table Grid'
        except Exception: pass
        set_table_borders(table_sig_1)
        
        s1_cells = table_sig_1.rows[0].cells
        
        # Left Cell formatting with bolding
        p_s1_left = s1_cells[0].paragraphs[0]
        p_s1_left.add_run("Prepared by:\n\n")
        p_s1_left.add_run("Enfrasys Solutions Sdn Bhd").bold = True
        p_s1_left.add_run("\n\nSignature:\n\n\n\n\n___________________________\nName:\nDesignation:\nDate:")
        
        # Right Cell formatting with bolding
        p_s1_right = s1_cells[1].paragraphs[0]
        p_s1_right.add_run("Verified by:\n\n")
        p_s1_right.add_run("Enfrasys Solutions Sdn Bhd").bold = True
        p_s1_right.add_run("\n\nSignature:\n\n\n\n\n___________________________\nName:\nDesignation:\nDate:")
        
        p_xml.addprevious(table_sig_1._element)
        
        # Page Break to jump to 2nd Signature Table
        p_break3 = add_p_before("")
        p_break3.add_run().add_break(WD_BREAK.PAGE)
        
        # TABLE 2 (Reviewed By / Verified / Approved)
        table_sig_2 = doc.add_table(rows=4, cols=2)
        try: table_sig_2.style = 'Table Grid'
        except Exception: pass
        set_table_borders(table_sig_2)
        
        r1_cells = table_sig_2.rows[0].cells
        r1_cells[0].merge(r1_cells[1])
        r1_cells[0].text = ""
        p_r1 = r1_cells[0].paragraphs[0]
        p_r1.add_run("Reviewed by:\n").bold = True
        p_r1.add_run("\n\n")
        
        r2_cells = table_sig_2.rows[1].cells
        r2_cells[0].text = ""
        r2_cells[0].paragraphs[0].add_run("Name").bold = True
        r2_cells[1].text = ""
        r2_cells[1].paragraphs[0].add_run("Signature").bold = True
        
        r3_cells = table_sig_2.rows[2].cells
        r3_cells[0].text = "Name:\n\nDesignation:\n\nDate:"
        r3_cells[1].text = "\n\n\n\n\n___________________________"
        
        r4_cells = table_sig_2.rows[3].cells
        r4_cells[0].text = ""
        p_r4_L = r4_cells[0].paragraphs[0]
        p_r4_L.add_run("Verified by:\n\n").bold = True
        p_r4_L.add_run("[Designation]\n\n[client]\n\nSignature:\n\n\n\n\n___________________________\nName:\n\nDesignation:\n\nDate:")
        
        r4_cells[1].text = ""
        p_r4_R = r4_cells[1].paragraphs[0]
        p_r4_R.add_run("Approved by:\n\n").bold = True
        p_r4_R.add_run("[Designation]\n\n[client]\n\nSignature:\n\n\n\n\n___________________________\nName:\n\nDesignation:\n\nDate:")
        
        p_xml.addprevious(table_sig_2._element)
        
        # Final Page Break before Executive Summary
        p_break4 = add_p_before("")
        p_break4.add_run().add_break(WD_BREAK.PAGE)
        
        # Now delete the original markdown '# Title' anchor so we don't have a blank Header 1 on Page 3
        p_elem = first_p._element
        p_elem.getparent().remove(p_elem)
        
        # Clean up the anchor we made for Page 2
        p_elem_anchor = anchor_p._element
        p_elem_anchor.getparent().remove(p_elem_anchor)
        
        # 1. THE ROBUST PAGE BREAK FIX & JUSTIFICATION
        print(f"  -> Applying page breaks & justifying text...")
        for para in doc.paragraphs:
            text = para.text.strip()
            
            # Justify all normal text
            if para.style.name == 'Normal' and text:
                para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            
            # Skip 1.0 from getting an automatic page break because our manual break pushes it to page 2 automatically
            if text.startswith("1.0"):
                continue
            
            # If paragraph is Heading 1 (e.g. 2.0, 3.0), force a page break! Sub-sections like 2.1 will remain on the same page.
            if para.style.name == 'Heading 1' or any(text.startswith(f"{i}.0") for i in range(2, 21)):
                para.paragraph_format.page_break_before = True
                
        # 2. THE TABLE FIX (Crash-Proof + Banding + SMART Column Widths!)
        print(f"  -> Fixing table layouts and applying smart column widths...")
        for table in doc.tables:
            # Check if this is a Signature table so we don't apply the blue header styling
            is_signature_table = False
            for row in table.rows:
                for cell in row.cells:
                    if "Prepared by:" in cell.text or "Reviewed by:" in cell.text:
                        is_signature_table = True
                        break
            
            if is_signature_table:
                try: table.style = 'Table Grid'
                except: pass
                continue

            try:
                table.style = 'Grid Table 4 - Accent 11' 
            except KeyError:
                try:
                    table.style = 'Grid Table 4 Accent 11'
                except KeyError:
                    table.style = 'Table Grid'
            
            # --- XML HACK 1: TURN OFF BANDED COLUMNS (Apply to ALL tables) ---
            tblPr = table._tbl.tblPr
            tblLook = tblPr.find(qn("w:tblLook"))
            if tblLook is None:
                tblLook = OxmlElement('w:tblLook')
                tblPr.append(tblLook)
                
            tblLook.set(qn('w:val'), '04A0')        
            tblLook.set(qn('w:firstRow'), '1')      
            tblLook.set(qn('w:lastRow'), '0')       
            tblLook.set(qn('w:firstColumn'), '1')   
            tblLook.set(qn('w:lastColumn'), '0')    
            tblLook.set(qn('w:noHBand'), '0')       
            tblLook.set(qn('w:noVBand'), '1')       
            
            # -----------------------------------------------------------------
            # NEW: ONLY STRETCH WIDE TABLES (5 or more columns)
            # -----------------------------------------------------------------
            if len(table.columns) >= 5:
                # --- XML HACK 2: FORCE WIDE TABLES TO 100% PAGE WIDTH ---
                tblW = tblPr.find(qn('w:tblW'))
                if tblW is None:
                    tblW = OxmlElement('w:tblW')
                    tblPr.append(tblW)
                tblW.set(qn('w:type'), 'pct')
                tblW.set(qn('w:w'), '5000')  # 100% width
                
                # --- XML HACK 3: DELETE PANDOC'S RIGID COLUMN WIDTHS ---
                for row in table.rows:
                    for cell in row.cells:
                        tcPr = cell._tc.get_or_add_tcPr()
                        tcW = tcPr.find(qn('w:tcW'))
                        if tcW is not None:
                            tcW.set(qn('w:type'), 'auto') 
                            tcW.set(qn('w:w'), '0')
            
            table.autofit = True
            table.allow_autofit = True

        # 3. THE SAFE DYNAMIC FOOTER TRICK
        print(f"  -> Applying dynamic footer title...")
        
        # 1. Clean the title to remove any invisible "Enter" keystrokes from the AI
        clean_title = document_title.replace('\n', '').replace('\r', '').strip()
        
        for section in doc.sections:
            for para in section.footer.paragraphs:
                for run in para.runs:
                    if '[PROJECT_TITLE]' in run.text:
                        run.text = run.text.replace('[PROJECT_TITLE]', clean_title)

        # 4. GLOBAL FONT ENFORCEMENT
        print(f"  -> Enforcing Segoe UI and typography scaling globally...")
        
        # Defined scale for headings (User specific)
        style_rules = {
            'Heading 1': {'size': 15, 'bold': False},
            'Heading 2': {'size': 14, 'bold': False}, 
            'Heading 3': {'size': 13, 'bold': False},
            'Heading 4': {'size': 12, 'bold': False},
        }

        # Format normal text, protect custom injected title pages!
        for para in doc.paragraphs:
            style_name = para.style.name
            rule = style_rules.get(style_name)

            for run in para.runs:
                run.font.name = 'Segoe UI' # Global override

                if rule: # It's a heading
                    run.font.size = Pt(rule['size'])
                    run.font.bold = rule['bold']
                else:
                    # It's Normal or custom. If it doesn't ALREADY have a hardcoded size (like our Title page has 22pt)
                    # then enforce 11pt exactly.
                    if run.font.size is None:
                        run.font.size = Pt(11)

        # Format tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        for run in p.runs:
                            run.font.name = 'Segoe UI'
                            # If no size explicitly set (like our custom Headers), force standard 11pt!
                            if run.font.size is None:
                                run.font.size = Pt(11)

        doc.save(str(docx_path))
        print(f"  ✓ Formatting successfully applied. Document saved to: {docx_path}")
        
    except ImportError as e:
        print(f"  ⚠ Missing library. Run: pip install pypandoc python-docx")
    except Exception as e:
        print(f"  ⚠ Failed to convert to Word: {e}")


# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────
def main():
    if len(sys.argv) > 1:
        pptx_path = sys.argv[1]
    else:
        # Combined script runs:
        import tkinter as tk
        from tkinter import filedialog
        
        root = tk.Tk()
        root.withdraw()
        root.attributes("-topmost", True)
        
        # 1. Ask for PPTX
        pptx_path = filedialog.askopenfilename(
            title="Select a PowerPoint file",
            initialdir=str(PROJECT_DIR),
            filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")],
        )
        
        # 2. Ask for Client Logo
        client_logo_path = None
        if pptx_path:
            client_logo_path = filedialog.askopenfilename(
                title="Select the Client Logo for Title Page (Optional)",
                initialdir=str(PROJECT_DIR),
                filetypes=[("Image files", "*.png *.jpg *.jpeg"), ("All files", "*.*")],
            )
            
        root.destroy()
        if not pptx_path:
            sys.exit(0)

    # If passed via command line, optionally accept a second argument as the client_logo_path
    if len(sys.argv) > 1:
        client_logo_path = sys.argv[2] if len(sys.argv) > 2 and sys.argv[2] else None

    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    print(f"=" * 60)
    print(f"  Phase 1: PPTX Content Extraction (via PDF & Vision)")
    print(f"  Input: {Path(pptx_path).name}")
    print(f"=" * 60)

    # --- NEW LOCAL POC CLEANUP ---
    if OUTPUT_DIR.exists():
        print(f"[CLEANUP] Deleting existing extraction-output folder...")
        shutil.rmtree(OUTPUT_DIR, ignore_errors=True)
    # -----------------------------

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    temp_dir = tempfile.mkdtemp(prefix="pptx_extract_")
    working_pptx_path = os.path.join(temp_dir, Path(pptx_path).name)
    working_pdf_path = os.path.join(temp_dir, Path(pptx_path).stem + ".pdf")
    
    shutil.copy2(pptx_path, working_pptx_path)

    try:
        # Step 1: python-pptx extraction
        print(f"\n[1/4] Extracting structure & images with python-pptx...")
        pptx_data = extract_with_pptx(working_pptx_path)
        print(f"  ✓ {pptx_data['metadata']['total_slides']} slides extracted")

        # Step 1.5: Convert to PDF
        convert_pptx_to_pdf(working_pptx_path, working_pdf_path)

        # Step 2: Content Understanding analysis (using the converted PDF)
        cu_data = analyze_with_content_understanding(working_pdf_path)

        if cu_data:
            contents = cu_data.get("result", cu_data).get("contents", [])
            if contents:
                md_path = OUTPUT_DIR / "content_understanding_output.md"
                with open(md_path, "w", encoding="utf-8") as f:
                    f.write(contents[0].get("markdown", ""))
                print(f"  ✓ Markdown output saved to {md_path}")

        # Step 2.5: Vision Analysis
        pptx_data["slides"] = analyze_images_with_vision(pptx_data["slides"])

        # Step 3: Merge results
        print(f"\n[4/4] Merging extraction results...")
        merged = merge_extraction_results(pptx_data, cu_data)

        merged_json_path = OUTPUT_DIR / "extraction_payload.json"
        with open(merged_json_path, "w", encoding="utf-8") as f:
            json.dump(merged, f, indent=2, ensure_ascii=False, default=str)
        print(f"  ✓ Final payload saved to {merged_json_path}")

        # ---------------------------------------------------------
        # NEW PHASE 2 & 3 PIPELINE
        # ---------------------------------------------------------
        
        # Step 4: Phase 2 (Orchestrator generates TOC)
        toc = generate_table_of_contents(merged)
        if toc:
            toc_path = OUTPUT_DIR / "table_of_contents.json"
            with open(toc_path, "w", encoding="utf-8") as f:
                json.dump(toc, f, indent=2)
            print(f"  ✓ Table of Contents saved to {toc_path}")

            # Step 5: Phase 3 (Writer drafts the document)
            final_doc_md = write_document_sections(toc, merged)
            doc_path = OUTPUT_DIR / "FINAL_DESIGN_DOCUMENT.md"
            with open(doc_path, "w", encoding="utf-8") as f:
                f.write(final_doc_md)
            print(f"  ✓ Final document generated and saved to {doc_path}")

            # Step 6: Phase 4 (Convert to Word)
            docx_path = OUTPUT_DIR / "Solution_Design_Document.docx"

            # Extract the dynamic client name and project title from the AI's outline
            client_name = toc.get("client_name", "").strip()
            project_title = toc.get("project_title", "Project").strip()
            
            # Combine them for the document title
            if client_name:
                doc_title = f"{client_name} {project_title}"
            else:
                doc_title = project_title

            # Pass the title to the conversion function
            convert_md_to_docx(doc_path, docx_path, doc_title, project_title, client_name, client_logo_path)

    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

    print(f"\n{'=' * 60}")
    print(f"  Pipeline Complete! Check your extraction-output folder.")
    print(f"{'=' * 60}\n")


if __name__ == "__main__":
    main()
