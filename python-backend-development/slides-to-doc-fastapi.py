import os
import sys
import json
import time
import base64
import requests
import shutil
import tempfile
import uuid
import re
import subprocess
from pathlib import Path
from threading import RLock
from typing import Optional, Dict, Any, List

from fastapi import FastAPI, UploadFile, File, BackgroundTasks, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pptx import Presentation
from dotenv import load_dotenv

try:
    from openai import AzureOpenAI  # noqa: F401
except ImportError:
    pass

load_dotenv()

# -----------------------------
# Environment & Paths
# -----------------------------
CONTENT_UNDERSTANDING_ENDPOINT = os.getenv("CONTENT_UNDERSTANDING_ENDPOINT", "").rstrip("/")
CONTENT_UNDERSTANDING_KEY = os.getenv("CONTENT_UNDERSTANDING_KEY", "")
API_VERSION = "2025-11-01"

AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT", "").rstrip("/")
AZURE_OPENAI_KEY = os.getenv("AZURE_OPENAI_KEY", "")
AZURE_OPENAI_DEPLOYMENT_NAME = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME", "gpt-4o")

AGENT_OPENAI_ENDPOINT = os.getenv("AGENT_OPENAI_ENDPOINT", "").rstrip("/")
AGENT_OPENAI_KEY = os.getenv("AGENT_OPENAI_KEY", "")

ORCHESTRATOR_DEPLOYMENT = os.getenv("ORCHESTRATOR_DEPLOYMENT", "gpt-4.1")
WRITER_DEPLOYMENT = os.getenv("WRITER_DEPLOYMENT", "gpt-4.1")
AGENT_ASSISTANT_ID = os.getenv("AGENT_ASSISTANT_ID", "")

# -----------------------------
# Pricing Details (USD) & Conversions
# -----------------------------
# Using standard pay-as-you-go Microsoft Foundry pricing for GPT-4.1
USD_TO_MYR_RATE = float(os.getenv("USD_TO_MYR_RATE", "4.75"))

# GPT-4.1 / GPT-4o Vision and Text pricing are currently unified under token count
RATE_VISION_PROMPT = 0.002 / 1000      # GPT-4.1 Global Text/Vision input ($2 per 1M)
RATE_VISION_COMPLETION = 0.008 / 1000  # GPT-4.1 Global Text/Vision output ($8 per 1M)
RATE_LLM_PROMPT = 0.002 / 1000         # GPT-4.1 Global Text input ($2 per 1M)
RATE_LLM_COMPLETION = 0.008 / 1000     # GPT-4.1 Global Text output ($8 per 1M)

# Content Understanding Base Extraction ($5 per 1k pages) + Estimated LLM Field Extraction
RATE_CU_PER_PAGE = 0.005               # Base extraction cost only

SCRIPT_DIR = Path(__file__).parent

# OUTPUT_DIR for Azure persistent storage
# - In Azure Web App Linux, set OUTPUT_DIR=/home/site/wwwroot/outputs and enable WEBSITES_ENABLE_APP_SERVICE_STORAGE=true
DEFAULT_OUTPUT_DIR = os.getenv("OUTPUT_DIR", str(SCRIPT_DIR / "outputs"))
BASE_STORAGE = Path(DEFAULT_OUTPUT_DIR)
BASE_STORAGE.mkdir(parents=True, exist_ok=True)

UPLOAD_DIR = BASE_STORAGE / "api_uploads"
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

TASKS_FILE = BASE_STORAGE / "tasks_db.json"
TASKS_LOCK = RLock()

# -----------------------------
# LibreOffice Path Detection
# -----------------------------
def _find_soffice() -> str:
    """Locate LibreOffice soffice binary. Works on both Linux and Windows."""
    # 1. Check if soffice is on PATH (Linux containers, or Windows with PATH configured)
    found = shutil.which("soffice")
    if found:
        return found

    # 2. Windows: check common installation paths
    if os.name == "nt":
        candidates = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        for candidate in candidates:
            if os.path.isfile(candidate):
                return candidate

    raise RuntimeError(
        "LibreOffice 'soffice' not found. "
        "Install LibreOffice and ensure it is on your PATH or in a standard location."
    )


# Cache the soffice path at module level (evaluated lazily on first call)
_SOFFICE_PATH: Optional[str] = None


def _get_soffice() -> str:
    """Return cached soffice path, finding it on first call."""
    global _SOFFICE_PATH
    if _SOFFICE_PATH is None:
        _SOFFICE_PATH = _find_soffice()
    return _SOFFICE_PATH


def _make_user_profile() -> str:
    """Create a temporary LibreOffice user profile directory.
    
    Each LibreOffice invocation needs its own profile to avoid
    lock conflicts when multiple requests run concurrently.
    """
    return tempfile.mkdtemp(prefix="lo_profile_")


# -----------------------------
# App & CORS
# -----------------------------
app = FastAPI(title="Design Document Generator API")

# Define the exact URLs that are allowed to talk to your backend.
# Do NOT put a trailing slash (/) at the end of the URL.
origins = [
    "http://localhost:3000",
    "https://design-document-converter.vercel.app"
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],  # Allows all HTTP methods (GET, POST, OPTIONS, etc.)
    allow_headers=["*"],  # Allows all headers (Content-Type, Authorization, etc.)
)

def _parse_cors_origins() -> List[str]:
    # Preserve your original default behavior (localhost only),
    # but allow Azure to configure via env: CORS_ORIGINS="https://yourfrontend,https://another"
    env_val = os.getenv("CORS_ORIGINS", "").strip()
    if not env_val:
        return ["http://localhost:3000"]
    return [o.strip() for o in env_val.split(",") if o.strip()]

def process_initial_upload(task_id: str, file_path: str, logo_path: Optional[str], src_name: str, task_dir: Path):
    try:
        update_task(task_id, {
            "status": "processing_upload",
            "step_name": "Generating Slide Thumbnails...",
            "progress": 20
        })

        # This is the slow part that was causing Vercel to timeout
        generate_slide_thumbnails(file_path, task_dir)

        update_task(task_id, {
            "step_name": "Extracting Slide Data...",
            "progress": 70
        })

        preview_data = extract_preview(file_path, task_id)

        # Mark as completely ready for the Next.js Preview Page
        update_task(task_id, {
            "status": "upload_complete",
            "step_name": "Ready for Preview",
            "progress": 100,
            "filename": src_name,
            "file_path": file_path,
            "logo_path": logo_path,
            "preview_data": preview_data
        })
    except Exception as e:
        update_task(task_id, {
            "status": "failed",
            "step_name": "Upload Failed",
            "error": str(e)
        })

app.add_middleware(
    CORSMiddleware,
    allow_origins=_parse_cors_origins(),
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# -----------------------------
# Startup diagnostic (kept as-is, safer behavior)
# -----------------------------
@app.on_event("startup")
async def test_agent_connection():
    print("\n" + "=" * 60)
    print("Initializing Application...")

    # Verify LibreOffice is available
    try:
        soffice = _get_soffice()
        print(f"[OK] LibreOffice found: {soffice}")
    except RuntimeError as e:
        print(f"[CRITICAL] {e}")

    if AGENT_ASSISTANT_ID:
        print(f"Detected AGENT_ASSISTANT_ID: {AGENT_ASSISTANT_ID}")
        print("Testing connection to Microsoft Foundry...")
        try:
            from azure.identity import DefaultAzureCredential
            from azure.ai.projects import AIProjectClient

            project_client = AIProjectClient(
                endpoint=AGENT_OPENAI_ENDPOINT,
                credential=DefaultAzureCredential()
            )
            # If this succeeds, credentials and endpoints are valid
            _openai_client = project_client.get_openai_client()

            agent_parts = AGENT_ASSISTANT_ID.split(":")
            ans_name = agent_parts[0]
            ans_ver = agent_parts[1] if len(agent_parts) > 1 else "2"

            print("[OK] FastAPI connected to Microsoft Foundry!")
            print(f" Bound to Agent: {ans_name} (v{ans_ver})")
            print(" (Bing Search will trigger automatically on generation)")
        except Exception as e:
            print("[WARNING] Failed to connect to Microsoft Foundry!")
            print(f" Reason: {str(e)}")
            print(" Fallback: Standard Offline GPT will be used.")
    else:
        print("[INFO] No AGENT_ASSISTANT_ID detected in .env.")
        print(" Fallback: Standard Offline GPT will be used.")
    print("=" * 60 + "\n")

@app.get("/")
def root_health_check():
    return {"status": "healthy", "message": "Enfrasys Document API is running"}

@app.get("/api/health")
async def api_health_check():
    return {"status": "healthy", "service": "Design Document Generator API"}


# -----------------------------
# Task Store (Azure-safe + less corruption risk)
# -----------------------------
def load_tasks() -> Dict[str, Any]:
    with TASKS_LOCK:
        if TASKS_FILE.exists():
            try:
                with open(TASKS_FILE, "r", encoding="utf-8") as f:
                    return json.load(f)
            except Exception:
                return {}
        return {}

def save_tasks(tasks_dict: Dict[str, Any]) -> None:
    with TASKS_LOCK:
        with open(TASKS_FILE, "w", encoding="utf-8") as f:
            json.dump(tasks_dict, f, indent=4, ensure_ascii=False, default=str)

def update_task(task_id: str, updates: Dict[str, Any]) -> None:
    with TASKS_LOCK:
        tasks = load_tasks()
        if task_id not in tasks:
            tasks[task_id] = {}
        tasks[task_id].update(updates)
        save_tasks(tasks)

def get_task(task_id: str) -> Optional[Dict[str, Any]]:
    tasks = load_tasks()
    return tasks.get(task_id)


# -----------------------------
# Filename sanitization (security hardening)
# -----------------------------
_FILENAME_SAFE_RE = re.compile(r"[^A-Za-z0-9._-]+")

def safe_filename(name: str, default: str = "file") -> str:
    # Keep base name only (prevents ../../ path traversal)
    base = Path(name).name.strip()
    if not base:
        base = default
    base = _FILENAME_SAFE_RE.sub("_", base)
    # Avoid extremely long names
    return base[:180]


# -----------------------------
# Conversions (LibreOffice-only, cross-platform)
# -----------------------------
def convert_pptx_to_pdf(pptx_path: str, pdf_path: str) -> None:
    """Convert PPTX to PDF using LibreOffice headless."""
    abs_pptx = os.path.abspath(pptx_path)
    abs_pdf = os.path.abspath(pdf_path)
    out_dir = os.path.dirname(abs_pdf)
    soffice = _get_soffice()
    user_profile = _make_user_profile()

    try:
        subprocess.run(
            [soffice, "--headless", "--nologo", "--nofirststartwizard",
             "--norestore",
             f"-env:UserInstallation=file:///{user_profile.replace(os.sep, '/')}",
             "--convert-to", "pdf", "--outdir", out_dir, abs_pptx],
            check=True, timeout=180,
            stdout=subprocess.PIPE, stderr=subprocess.PIPE
        )
        expected_out = os.path.join(out_dir, os.path.splitext(os.path.basename(abs_pptx))[0] + ".pdf")
        if expected_out != abs_pdf and os.path.exists(expected_out):
            os.rename(expected_out, abs_pdf)
        if not os.path.exists(abs_pdf):
            raise RuntimeError("LibreOffice conversion succeeded but output PDF not found.")
    except subprocess.TimeoutExpired:
        raise RuntimeError("PPTX to PDF conversion timed out (180s).")
    except Exception as e:
        raise RuntimeError(f"PPTX to PDF conversion failed: {e}")
    finally:
        shutil.rmtree(user_profile, ignore_errors=True)


def convert_docx_to_pdf(docx_path: str, pdf_path: str) -> None:
    """Convert DOCX to PDF using LibreOffice headless."""
    abs_docx = os.path.abspath(docx_path)
    abs_pdf = os.path.abspath(pdf_path)
    out_dir = os.path.dirname(abs_pdf)
    soffice = _get_soffice()
    user_profile = _make_user_profile()

    try:
        subprocess.run(
            [soffice, "--headless", "--norestore", "--nologo",
             f"-env:UserInstallation=file:///{user_profile.replace(os.sep, '/')}",
             "--convert-to", "pdf",
             "--outdir", out_dir, abs_docx],
            check=True, timeout=180,
            stdout=subprocess.PIPE, stderr=subprocess.PIPE
        )

        expected_out = os.path.join(out_dir, os.path.splitext(os.path.basename(abs_docx))[0] + ".pdf")
        if expected_out != abs_pdf and os.path.exists(expected_out):
            os.rename(expected_out, abs_pdf)

        if not os.path.exists(abs_pdf):
            raise RuntimeError("LibreOffice conversion succeeded but output PDF not found.")
    except subprocess.TimeoutExpired:
        raise RuntimeError("DOCX to PDF conversion timed out (180s).")
    except Exception as e:
        raise RuntimeError(f"DOCX to PDF conversion failed: {e}")
    finally:
        shutil.rmtree(user_profile, ignore_errors=True)


def _update_toc_with_libreoffice(docx_path: str) -> None:
    """Fallback TOC update: round-trip DOCX through LibreOffice to force field recalculation."""
    abs_path = os.path.abspath(docx_path)
    out_dir = os.path.dirname(abs_path)
    base_name = os.path.splitext(os.path.basename(abs_path))[0]
    soffice = _get_soffice()

    # Create a backup
    backup_path = abs_path + ".bak"
    shutil.copy2(abs_path, backup_path)

    try:
        # Convert DOCX -> DOCX through LibreOffice (forces field update)
        new_profile = _make_user_profile()
        try:
            subprocess.run(
                [soffice, "--headless", "--norestore",
                 f"-env:UserInstallation=file:///{new_profile.replace(os.sep, '/')}",
                 "--infilter=Microsoft Word 2007-2019 XML",
                 "--convert-to", "docx:Microsoft Word 2007-2019 XML",
                 "--outdir", out_dir, abs_path],
                check=True, timeout=120,
                stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )
        finally:
            shutil.rmtree(new_profile, ignore_errors=True)

        # LibreOffice may produce the output with the same name, check if it exists
        if os.path.exists(abs_path):
            # Remove backup since conversion succeeded
            os.remove(backup_path)
        else:
            # Restore from backup
            os.rename(backup_path, abs_path)
    except Exception as e:
        # Restore from backup on failure
        if os.path.exists(backup_path):
            if os.path.exists(abs_path):
                os.remove(abs_path)
            os.rename(backup_path, abs_path)
        print(f"[TOC Roundtrip] Warning: {e}")
        raise e


# -----------------------------
# Thumbnails / Preview
# -----------------------------
def generate_slide_thumbnails(pptx_path: str, output_dir: Path) -> None:
    import fitz

    temp_pdf = str(output_dir / "temp_thumbnail_source.pdf")
    convert_pptx_to_pdf(pptx_path, temp_pdf)

    if not os.path.exists(temp_pdf):
        print("Failed to generate PDF for thumbnail extraction.")
        return

    doc = None
    try:
        doc = fitz.open(temp_pdf)
        for i in range(len(doc)):
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            export_path = str(output_dir / f"slide_{i+1}.jpg")
            pix.save(export_path)
    except Exception as e:
        print(f"Failed to generate thumbnails via PyMuPDF: {e}")
    finally:
        try:
            if doc is not None:
                doc.close()
        except Exception:
            pass
        try:
            if os.path.exists(temp_pdf):
                os.remove(temp_pdf)
        except Exception:
            pass


def extract_preview(pptx_path: str, task_id: str) -> Dict[str, Any]:
    prs = Presentation(pptx_path)
    slides_data = []
    diagram_count = 0

    for idx, slide in enumerate(prs.slides, 1):
        title = ""
        notes = "No presenter note"

        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
            n = slide.notes_slide.notes_text_frame.text.strip()
            if n:
                notes = n

        for shape in slide.shapes:
            if shape.shape_type == 13:
                diagram_count += 1

        slides_data.append({
            "slide_number": idx,
            "title": title if title else f"Slide {idx}",
            "notes": notes,
            "thumbnail_url": f"/api/thumbnails/{task_id}/slide_{idx}.jpg"
        })

    return {
        "filename": Path(pptx_path).name,
        "total_slides": len(prs.slides),
        "diagrams_detected": diagram_count,
        "slides": slides_data
    }


# -----------------------------
# PPTX extraction
# -----------------------------
def extract_shapes(shapes, slide_info: Dict[str, Any], ctx: Dict[str, Any]) -> None:
    for shape in shapes:
        # GroupShape recursion
        if shape.shape_type == 6:
            extract_shapes(shape.shapes, slide_info, ctx)
            continue

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text and text != slide_info.get("title"):
                    slide_info["text_content"].append(text)

        if shape.shape_type == 13:
            ctx["image_counter"] += 1
            img = shape.image
            ext = img.content_type.split("/")[-1]
            if ext in ["x-wmf", "wmf", "x-emf", "emf"]:
                ext = "png"

            filename = f"slide_{ctx['slide_idx']:03d}_img_{ctx['image_counter']:03d}.{ext}"
            filepath = ctx["images_dir"] / filename

            if ext in ["png"]:
                try:
                    import io
                    from PIL import Image
                    wmf_img = Image.open(io.BytesIO(img.blob))
                    wmf_img.save(filepath, format="PNG")
                    content_type = "image/png"
                    size_bytes = os.path.getsize(filepath)
                except Exception:
                    filename = f"slide_{ctx['slide_idx']:03d}_img_{ctx['image_counter']:03d}.x-wmf"
                    filepath = ctx["images_dir"] / filename
                    with open(filepath, "wb") as f:
                        f.write(img.blob)
                    content_type = img.content_type
                    size_bytes = len(img.blob)
            else:
                if ext == "jpeg":
                    ext = "jpg"
                filename = f"slide_{ctx['slide_idx']:03d}_img_{ctx['image_counter']:03d}.{ext}"
                filepath = ctx["images_dir"] / filename
                with open(filepath, "wb") as f:
                    f.write(img.blob)
                content_type = img.content_type
                size_bytes = len(img.blob)

            slide_info["images"].append({
                "image_id": f"img_{ctx['slide_idx']:03d}_{ctx['image_counter']:03d}",
                "filename": filename,
                "file_path": str(filepath),
                "content_type": content_type,
                "size_bytes": size_bytes,
            })

        if shape.has_table:
            table = shape.table
            headers = [cell.text.strip() for cell in table.rows[0].cells] if table.rows else []
            rows = []
            for row_idx, row in enumerate(table.rows):
                if row_idx == 0:
                    continue
                row_data = [cell.text.strip() for cell in row.cells]
                rows.append(row_data)

            slide_info["tables"].append({
                "headers": headers,
                "rows": rows,
                "total_rows": len(table.rows),
                "total_cols": len(table.columns),
            })


def extract_with_pptx(pptx_path: str, output_dir: Path) -> Dict[str, Any]:
    prs = Presentation(pptx_path)
    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    slides_data = []
    ctx = {"image_counter": 0, "images_dir": images_dir, "slide_idx": 0}

    for slide_idx, slide in enumerate(prs.slides, 1):
        ctx["slide_idx"] = slide_idx
        slide_info = {
            "slide_number": slide_idx,
            "title": None,
            "text_content": [],
            "speaker_notes": None,
            "images": [],
            "tables": []
        }

        if slide.shapes.title:
            slide_info["title"] = slide.shapes.title.text.strip()

        extract_shapes(slide.shapes, slide_info, ctx)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_info["speaker_notes"] = notes

        slides_data.append(slide_info)

    return {
        "metadata": {
            "source_file": Path(pptx_path).name,
            "total_slides": len(prs.slides),
            "total_images_extracted": ctx["image_counter"],
            "slides_with_notes": sum(1 for s in slides_data if s["speaker_notes"]),
        },
        "slides": slides_data,
    }


# -----------------------------
# Content Understanding
# -----------------------------
def analyze_with_content_understanding(pdf_path: str) -> Optional[Dict[str, Any]]:
    if not CONTENT_UNDERSTANDING_ENDPOINT or not CONTENT_UNDERSTANDING_KEY:
        return None

    analyze_binary_url = (
        f"{CONTENT_UNDERSTANDING_ENDPOINT}/contentunderstanding/analyzers/"
        f"design_document_converter:analyzeBinary?api-version={API_VERSION}"
    )

    with open(pdf_path, "rb") as f:
        file_bytes = f.read()

    headers = {
        "Ocp-Apim-Subscription-Key": CONTENT_UNDERSTANDING_KEY,
        "Content-Type": "application/pdf",
    }

    response = requests.post(analyze_binary_url, headers=headers, data=file_bytes)
    if response.status_code not in (200, 202):
        return _analyze_with_base64_fallback(pdf_path)

    operation_url = response.headers.get("Operation-Location")
    if not operation_url:
        if response.status_code == 200:
            return response.json()
        return None

    return _poll_for_result(operation_url)


def _analyze_with_base64_fallback(pdf_path: str) -> Optional[Dict[str, Any]]:
    analyze_url = (
        f"{CONTENT_UNDERSTANDING_ENDPOINT}/contentunderstanding/analyzers/"
        f"design_document_converter:analyze?api-version={API_VERSION}"
    )

    with open(pdf_path, "rb") as f:
        file_bytes = f.read()

    b64_content = base64.b64encode(file_bytes).decode("utf-8")
    data_uri = f"data:application/pdf;base64,{b64_content}"

    headers = {
        "Ocp-Apim-Subscription-Key": CONTENT_UNDERSTANDING_KEY,
        "Content-Type": "application/json",
    }

    payload = {"inputs": [{"data": data_uri}]}
    response = requests.post(analyze_url, headers=headers, json=payload)

    if response.status_code not in (200, 202):
        return None

    operation_url = response.headers.get("Operation-Location")
    if operation_url:
        return _poll_for_result(operation_url)

    return response.json() if response.status_code == 200 else None


def _poll_for_result(operation_url: str, max_retries: int = 60, interval: int = 3) -> Optional[Dict[str, Any]]:
    headers = {"Ocp-Apim-Subscription-Key": CONTENT_UNDERSTANDING_KEY}
    for _ in range(max_retries):
        time.sleep(interval)
        response = requests.get(operation_url, headers=headers)
        if response.status_code != 200:
            continue
        result = response.json()
        status = result.get("status", "")
        if status == "Succeeded":
            return result
        if status in ("Failed", "Cancelled"):
            return None
    return None


# -----------------------------
# Vision Image Analysis
# -----------------------------
def analyze_images_with_vision(task_id: str, slides_data: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    if not AZURE_OPENAI_ENDPOINT or not AZURE_OPENAI_KEY:
        return slides_data

    base_endpoint = AZURE_OPENAI_ENDPOINT.split("/openai")[0].rstrip("/")
    url = f"{base_endpoint}/openai/v1/chat/completions"
    headers = {
        "api-key": AZURE_OPENAI_KEY,
        "Authorization": f"Bearer {AZURE_OPENAI_KEY}",
        "Content-Type": "application/json"
    }

    vision_prompt = (
        "You are an Expert Azure Solutions Architect. Analyze this image extracted from a pre-sales presentation.\n"
        "1. If this image is a company logo, a decorative background, a transition slide, or a generic stock photo, reply ONLY with the word: DECORATIVE.\n"
        "2. If this image is a technical architecture diagram, a data workflow, or a system component map, provide a highly detailed, professional text description of the data flow, the components involved, and how they connect.\n"
        "Do not include conversational filler."
    )

    for slide in slides_data:
        for img in slide.get("images", []):
            filepath = img.get("file_path")
            if not filepath or not os.path.exists(filepath):
                continue

            with open(filepath, "rb") as f:
                base64_image = base64.b64encode(f.read()).decode("utf-8")

            mime_type = img.get("content_type", "image/png")

            payload = {
                "model": AZURE_OPENAI_DEPLOYMENT_NAME,
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": vision_prompt},
                            {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{base64_image}"}}
                        ]
                    }
                ],
                "max_tokens": 800,
                "temperature": 0.0
            }

            try:
                response = requests.post(url, headers=headers, json=payload, timeout=120)
                if response.status_code == 200:
                    data = response.json()
                    description = data["choices"][0]["message"]["content"].strip()
                    img["ai_description"] = description
                    
                    # Accumulate token usage
                    usage = data.get("usage", {})
                    task = get_task(task_id)
                    if task and "cost_metrics" in task:
                        task["cost_metrics"]["vision_tokens_prompt"] += usage.get("prompt_tokens", 0)
                        task["cost_metrics"]["vision_tokens_completion"] += usage.get("completion_tokens", 0)
                        
            except Exception:
                # preserve original behavior (silent on error)
                pass

    return slides_data


def merge_extraction_results(pptx_data: Dict[str, Any], cu_data: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    result = {
        "metadata": pptx_data["metadata"],
        "extraction_sources": ["python-pptx", "azure-openai-vision"],
        "slides": pptx_data["slides"],
        "content_understanding": None,
    }

    if cu_data:
        result["extraction_sources"].append("content-understanding")
        cu_result = cu_data.get("result", cu_data)
        contents = cu_result.get("contents", [])
        if contents:
            markdown_content = contents[0].get("markdown", "")
            result["content_understanding"] = {
                "markdown": markdown_content,
                "markdown_length": len(markdown_content),
            }

    return result


# -----------------------------
# Orchestrator & Writer
# -----------------------------
def generate_table_of_contents(task_id: str, extraction_payload: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    if not AGENT_OPENAI_ENDPOINT or not AGENT_OPENAI_KEY:
        return None

    base_endpoint = AGENT_OPENAI_ENDPOINT.split("/openai")[0].rstrip("/")
    url = f"{base_endpoint}/openai/v1/chat/completions"
    headers = {
        "api-key": AGENT_OPENAI_KEY,
        "Authorization": f"Bearer {AGENT_OPENAI_KEY}",
        "Content-Type": "application/json"
    }

    orchestrator_prompt = (
        "You are an Expert Enterprise Solutions Architect from Enfrasys Solutions. Your task is to analyze an extracted solution design presentation "
        "(provided as a combined JSON payload) and produce a detailed Table of Contents (Outline) for a formal Microsoft Word Solution Design Document that follows the Enfrasys standard structure.\n\n"
        "INSTRUCTIONS:\n"
        "1. Analyze the payload deeply: read slide text, speaker notes, and diagram descriptions to fully understand the technical context and design decisions made.\n"
        "2. Group related slides into the standard Enfrasys architectural sections listed below. Do NOT map slides 1-to-1.\n"
        "3. REQUIRED SECTION ORDER — Follow this exact pattern, including strict decimal numbering (e.g., 2.1, 2.2) for all sub-sections. Adapt section titles to the technology in the presentation:\n"
        " 2.0 Executive Summary\n"
        " Sub-sections: 2.1 Project Overview 2.2 Document Purpose 2.3 Document Audience\n"
        " 3.0 [Technology] Overview\n"
        " Sub-sections: 3.1 [Core Technology 1] 3.2 [Core Technology 2] (Force the Writer to generate boilerplate definitions for the core Microsoft technologies used, e.g., Fabric, M365, Entra ID).\n"
        " 4.0 Network Design and Decision\n"
        " Sub-sections: 4.1 Network Connectivity Overview 4.2 [Technology] Data Gateway (with VM spec table if a gateway VM exists)\n"
        " 5.0 Roles in [Technology]\n"
        " Sub-sections: 5.1 [Technology] Administrator Role 5.2 [Technology]-native Roles 5.3 [Client] Workspace Role and Access\n"
        " 6.0 [Technology] Design and Decision ← This is the LARGEST section. It MUST cover ALL of the following sub-sections:\n"
        " 6.1 [Topic] Workflow Design Considerations (table of design considerations with IDs and descriptions)\n"
        " 6.2 [Topic] Workflow Design Decisions\n"
        " 6.3 Access Design Considerations and Decisions (H3 per access method, plus decision table)\n"
        " 6.4 Virtual Machine (VM) Size (table with OS, .NET, vCPU, RAM specs for any gateway/VM component)\n"
        " 6.5 Resource Organization Design (H3: Management Groups H3: Naming and Tagging → tables: Azure Naming, [Tech] Naming, Acronyms)\n"
        " 6.6 Governance Consideration (H3: Governance Disciplines with Cost Management table H3: Violation Triggers and Actions table)\n"
        " 7.0 Security Design and Decision\n"
        " Sub-sections: 7.1 Azure NSG Overview (Inbound/Outbound Rules tables) 7.2 Encryption Design\n"
        " 8.0 Deployment & Migration Approach [If Applicable]\n"
        " Sub-sections: 8.1 Pre-Migration/Setup 8.2 Pilot 8.3 Production 8.4 Post-Migration/Support. (Instruct the Writer to generate 'Action By' task tables and D-7/D0/D+1 timelines here).\n"
        " 9.0 Appendix\n"
        " Sub-sections: 9.1 Appendix 1 Computing 9.2 Appendix 2 Network 9.3 Appendix 3 Identity & Security 9.4 Appendix 4 Logging & Monitoring 9.5 Appendix 5 Cloud Governance 9.6 Appendix 6 [Tech-specific]\n"
        " EACH appendix sub-section MUST have: 'Introduction/Prerequisites', 'Limits and Boundaries', 'Others'.\n"
        "4. CONTENT PLACEMENT RULES:\n"
        " - NSG rules belong in Security Design ONLY.\n"
        " - VM specifications belong in the Design and Decision section ONLY.\n"
        " - Naming conventions belong in Resource Organization Design ONLY.\n"
        "5. Do NOT create standalone top-level sections for 'Microsoft Best Practices', 'Governance', or 'Naming Conventions'.\n"
        "6. Do NOT generate Section 1.0 (Document Sign Off). Start ALL section numbering at 2.0.\n"
        "7. Output the outline strictly as a JSON object matching this schema:\n"
        "{\n"
        ' "client_name": "string (Exact client name from title slide)",\n'
        ' "client_name_full": "string | null",\n'
        ' "project_title": "string",\n'
        ' "sections": [\n'
        "  {\n"
        '   "section_number": "2.0",\n'
        '   "section_title": "string",\n'
        '   "mapped_slides": [int],\n'
        '   "generation_instructions": "string (List every H2/H3 sub-section with its exact decimal number (e.g., 2.1, 6.3), required tables WITH column names, explicit instructions to add boilerplate definitions, and instructions to assign Enfrasys vs. Client responsibilities.)"\n'
        "  }\n"
        " ]\n"
        "}\n"
    )

    payload = {
        "model": ORCHESTRATOR_DEPLOYMENT,
        "response_format": {"type": "json_object"},
        "messages": [
            {"role": "system", "content": orchestrator_prompt},
            {"role": "user", "content": json.dumps(extraction_payload, default=str)}
        ],
        "temperature": 0.2
    }

    try:
        response = requests.post(url, headers=headers, json=payload, timeout=180)
        if response.status_code == 200:
            data = response.json()
            content = data["choices"][0]["message"]["content"]
            
            # Accumulate token usage
            usage = data.get("usage", {})
            task = get_task(task_id)
            if task and "cost_metrics" in task:
                task["cost_metrics"]["llm_tokens_prompt"] += usage.get("prompt_tokens", 0)
                task["cost_metrics"]["llm_tokens_completion"] += usage.get("completion_tokens", 0)
                
            return json.loads(content)
        return None
    except Exception:
        return None


def write_document_sections(task_id: str, toc: Dict[str, Any], extraction_payload: Dict[str, Any]) -> str:
    base_endpoint = AGENT_OPENAI_ENDPOINT.split("/openai")[0].rstrip("/")
    url = f"{base_endpoint}/openai/v1/chat/completions"
    headers = {
        "api-key": AGENT_OPENAI_KEY,
        "Authorization": f"Bearer {AGENT_OPENAI_KEY}",
        "Content-Type": "application/json"
    }

    final_document_markdown = f"# {toc.get('client_name', '')} {toc.get('project_title', 'Solution Design Document')}".strip() + "\n\n"

    writer_system_prompt = """
        You are a Lead Enterprise Architect from Enfrasys Solutions writing a formal Solution Design Document. Your output will be converted directly to a Microsoft Word document for a client.
        Write in authoritative, professional language ("Enfrasys recommends", "[CLIENT_NAME] shall implement"). Be specific — reference the actual client name, technology components, and design decisions from the slide data.
        CRITICAL: Do NOT wrap your entire response in a ```markdown code block. Just return the raw text.
        
        -- CONSULTING EXPANSION RULES (CRITICAL) --
        1. STATIC BOILERPLATE: Whenever you introduce a major Microsoft technology (e.g., Microsoft Fabric, Exchange Online, Entra ID, Power BI), you MUST provide a standard, formal definition of that technology before discussing the client's specific design. Do not assume the reader knows what the technology does.
        2. RESPONSIBILITY ASSIGNMENT: In any section discussing tasks, migrations, or rollouts, you MUST explicitly assign ownership. Use "Enfrasys Solutions" for vendor tasks and "[CLIENT_NAME]" for client tasks.
        3. TIMELINES (If Applicable): If the slides mention migration waves or rollouts, automatically structure them into formal consulting timelines (e.g., Pre-Migration (D-7), Cutover (D0), Post-Migration Support (D+1)).

        -- MARKDOWN HEADING RULES --
        - Main Sections MUST use `#` (e.g., `# 5.0 Fabric Design and Decision`).
        - Sub-sections MUST use `##` (e.g., `## 5.1 Data Workflow Design Considerations`).
        - Sub-sub-sections MUST use `###` (e.g., `### 5.7.1 Azure Management Group`).
        - The H1 main section opening paragraph must be 1-2 sentences only. Detail goes in H2/H3 sub-sections.

        -- MANDATORY EXECUTION WORKFLOW --
        You possess the Web Search tool. You MUST execute your task in the following strict sequence. Do not skip any steps.
        1. SEARCH: You must use your Web Search tool to query MS Learn and official Microsoft forums for the latest limits, capabilities, and naming conventions for the technology mentioned. Do NOT rely entirely on your internal knowledge.
        2. DRAFT: Write the document section based on both the slide data and the live search results.
        3. RECHECK: Before finalizing, evaluate your draft internally. Are the technical specifications up to date based on your search? Are you following the Markdown and Table constraints?
        4. OUTPUT: Proceed to output the final, corrected text.

        -- TABLE RULES (STRICT SCHEMA) --
        Use Markdown tables to present all technical data. You must use the exact columns below when generating these tables:
        - Design Considerations: columns = ID | Description | Workload Type
        - Design Decisions: columns = No. | Design Decision | Decision
        - Task/Migration Rollout: columns = Item | Activities | Action By | Status (Assign 'Action By' to Enfrasys or [CLIENT_NAME])
        - Administrator capabilities: columns = Capability | Description
        - Workspace-level role capabilities: columns = Permission | Admin | Contributor | Member | Viewer
        - Client workspace role mapping: columns = No. | Workspace Role | Suggested Role | [CLIENT_NAME] Personnel
        - Access decisions: columns = No. | Policy | Decision
        - VM specifications: columns = Component | Specification
        - NSG rules: columns = Name | Priority | Source | Source Ports | Destination | Destination Ports | Protocol | Access
        - Azure Naming Convention: columns = Resource | Abbreviation | Example
        - Cost Management Tools: columns = Azure Tool | Description | Cost Management Discipline

        -- STRICT IMAGE RULES --
        - Select only architecture diagrams, data flow diagrams, network topologies, or access/security diagrams.
        - Each UNIQUE image file_path may only be embedded ONCE across the entire document.
        - Syntax: `![](file_path)` — use the exact path from the JSON. Do NOT invent paths.
        - CAPTION RULE: You MUST place a blank empty line between the image and its caption to separate them! 
            Example:
            ![](file_path)

            Figure 1: High-level architecture.
        - DO NOT copy-paste the long 'ai_description' from the JSON into the document. Keep the actual output caption extremely brief.

        -- NO-REPETITION RULES --
        1. Each table must appear EXACTLY ONCE across the full document. If a table was in a main section, the Appendix must NOT repeat it.
        2. NSG rules belong ONLY in the Security Design section.
        3. VM specifications belong ONLY in the Platform Design section.
        4. Naming convention tables belong ONLY in the Resource Organization Design sub-section.
        
        -- APPENDIX FORMATTING (STRICT RULES) --
        Each appendix entry MUST have exactly these three H3 sub-sections:
        ### Introduction/Prerequisites
        ### Limits and Boundaries
        ### Others
        ZERO EXPLANATION RULE: Do NOT write any introductory sentences, concluding remarks, or explanations in the Appendix sections. Output ONLY the headers and the links.
        
        CRITICAL LINK FORMATTING (READ CAREFULLY):
        - You MUST leave a blank empty line between every single link so they do not merge together into one paragraph.
        - Do NOT use bullet points (`*` or `-`).
        - To ensure the URL is explicitly visible to the reader AND acts as a clickable hyperlink in the Word document, you MUST use the Markdown link syntax where the URL is BOTH the display text and the link destination.
        - Use this EXACT format: Title of the Reference: [https://...](https://...)
        
        Example:
        Azure Compute Overview: [https://learn.microsoft.com/en-us/azure/virtual-machines/](https://learn.microsoft.com/en-us/azure/virtual-machines/)
    """

    use_agent = False
    openai_client = None
    agent_name = ""
    agent_version = "2"

    if AGENT_ASSISTANT_ID:
        try:
            from azure.identity import DefaultAzureCredential
            from azure.ai.projects import AIProjectClient

            project_client = AIProjectClient(
                endpoint=AGENT_OPENAI_ENDPOINT,
                credential=DefaultAzureCredential()
            )
            openai_client = project_client.get_openai_client()
            agent_parts = AGENT_ASSISTANT_ID.split(":")
            agent_name = agent_parts[0]
            agent_version = agent_parts[1] if len(agent_parts) > 1 else "2"
            use_agent = True
        except Exception as e:
            print(f"Failed to intialize Azure AI Projects Client: {e}")
            use_agent = False

    for section in toc.get("sections", []):
        sec_num = section.get("section_number")
        sec_title = section.get("section_title")
        instructions = section.get("generation_instructions")
        mapped_slides = section.get("mapped_slides", [])

        relevant_slide_data = [
            s for s in extraction_payload.get("slides", [])
            if s.get("slide_number") in mapped_slides
        ]

        client_name = toc.get("client_name", "")
        client_name_full = toc.get("client_name_full", "")

        context_str = f"Client Acronym: {client_name}\n"
        if client_name_full:
            context_str += f"Client Full Name: {client_name_full} (Use this full name in the content when appropriate)\n"

        user_prompt = (
            f"Write section '{sec_num} {sec_title}' for the design document.\n"
            f"{context_str}\n"
            f"Instructions from Orchestrator: {instructions}\n\n"
            f"Here is the raw slide data (text, speaker notes, and diagram descriptions):\n"
            f"{json.dumps(relevant_slide_data, default=str)}"
        )

        drafted_text = None

        if use_agent and openai_client:
            try:
                response = openai_client.responses.create(
                    input=[
                        {"role": "system", "content": writer_system_prompt},
                        {"role": "user", "content": user_prompt}
                    ],
                    extra_body={
                        "agent_reference": {
                            "name": agent_name,
                            "version": agent_version,
                            "type": "agent_reference"
                        }
                    }
                )
                drafted_text = response.output_text
            except Exception:
                pass

        if not use_agent or not drafted_text:
            payload = {
                "model": WRITER_DEPLOYMENT,
                "messages": [
                    {"role": "system", "content": writer_system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                "temperature": 0.4
            }
            try:
                response = requests.post(url, headers=headers, json=payload, timeout=240)
                if response.status_code == 200:
                    data = response.json()
                    drafted_text = data["choices"][0]["message"]["content"].strip()
                    
                    # Accumulate token usage (fallback API)
                    usage = data.get("usage", {})
                    task = get_task(task_id)
                    if task and "cost_metrics" in task:
                        task["cost_metrics"]["llm_tokens_prompt"] += usage.get("prompt_tokens", 0)
                        task["cost_metrics"]["llm_tokens_completion"] += usage.get("completion_tokens", 0)
            except Exception:
                pass

        if drafted_text:
            if drafted_text.startswith("```"):
                drafted_text = drafted_text.split("\n", 1)[-1]
            if drafted_text.endswith("```"):
                drafted_text = drafted_text.rsplit("\n", 1)[0]
            drafted_text = drafted_text.strip()
            final_document_markdown += f"{drafted_text}\n\n"
        else:
            final_document_markdown += f"## {sec_num} {sec_title}\n\n*(Error generating section)*\n\n"

    return final_document_markdown


# =============================================================================
# Markdown -> DOCX (Pandoc + python-docx post-processing)
# =============================================================================
# This section is structured as a pipeline of focused sub-functions for
# maintainability and testability.

def _pandoc_md_to_raw_docx(md_path: str, docx_path: str) -> None:
    """Phase 1: Use Pandoc to convert Markdown to a raw DOCX with the Enfrasys template."""
    import pypandoc

    template_file = str(SCRIPT_DIR / "enfrasys_template.docx")
    if not os.path.exists(template_file):
        raise FileNotFoundError(f"Deployment Error: enfrasys_template.docx is missing at {template_file}")

    task_dir_str = str(Path(md_path).parent.resolve())
    extra_args = [
        f"--reference-doc={template_file}",
        f"--resource-path={task_dir_str}"
    ]

    pypandoc.convert_file(
        str(md_path),
        "docx",
        outputfile=str(docx_path),
        extra_args=extra_args
    )


def _build_title_page(doc, first_p, project_title, document_title, client_name, client_logo_path):
    """Build the cover/title page at the very front of the document."""
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    first_p.paragraph_format.page_break_before = False

    def set_p_bottom_border(p, sz="144"):
        pPr = p._element.get_or_add_pPr()
        pBdr = OxmlElement("w:pBdr")
        pPr.append(pBdr)
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), sz)
        bottom.set(qn("w:space"), "1")
        bottom.set(qn("w:color"), "000000")
        pBdr.append(bottom)

    # Thick bar
    p_bar = first_p.insert_paragraph_before("")
    set_p_bottom_border(p_bar)

    # Project title
    p_title = first_p.insert_paragraph_before()
    p_title.paragraph_format.space_before = Pt(6)
    run = p_title.add_run(project_title)
    run.font.name = "Segoe UI"
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(0, 0, 0)

    # Subtitle
    p_sub = first_p.insert_paragraph_before()
    p_sub.paragraph_format.space_before = Pt(24)
    run = p_sub.add_run("Planning & Design Document")
    run.font.name = "Segoe UI"
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(0, 0, 0)

    # Client logo
    if client_logo_path and os.path.exists(client_logo_path):
        p_logo = first_p.insert_paragraph_before()
        p_logo.alignment = WD_ALIGN_PARAGRAPH.LEFT
        p_logo.paragraph_format.space_before = Pt(40)
        run = p_logo.add_run()
        try:
            run.add_picture(str(client_logo_path), height=Inches(1.5))
        except Exception:
            pass

    # Spacing
    for _ in range(4):
        first_p.insert_paragraph_before("")

    # Thin divider
    p_abs_div = first_p.insert_paragraph_before("")
    set_p_bottom_border(p_abs_div, sz="8")

    # Abstract
    p_abs_title = first_p.insert_paragraph_before()
    p_abs_title.paragraph_format.space_before = Pt(6)
    run = p_abs_title.add_run("Abstract")
    run.font.name = "Segoe UI"
    run.bold = True
    run.font.size = Pt(11)

    p_abs_text = first_p.insert_paragraph_before()
    p_abs_text.paragraph_format.space_before = Pt(6)
    run = p_abs_text.add_run(f"This document defines the Planning and Design for {document_title} project.")
    run.font.name = "Segoe UI"
    run.font.size = Pt(11)

    for _ in range(2):
        first_p.insert_paragraph_before("")

    # Prepared By
    p_prep = first_p.insert_paragraph_before()
    run = p_prep.add_run("Prepared By")
    run.font.name = "Segoe UI"
    run.font.size = Pt(11)

    enfrasys_logo = SCRIPT_DIR / "enfrasys-logo.jpg"
    if os.path.exists(enfrasys_logo):
        p_enfrasys = first_p.insert_paragraph_before()
        p_enfrasys.alignment = WD_ALIGN_PARAGRAPH.LEFT
        p_enfrasys.paragraph_format.space_before = Pt(10)
        run = p_enfrasys.add_run()
        try:
            run.add_picture(str(enfrasys_logo), width=Inches(1.5))
        except Exception:
            pass

    # Page break after title page
    p_break_title = first_p.insert_paragraph_before("")
    p_break_title.add_run().add_break(WD_BREAK.PAGE)


def _apply_header_footer(doc, document_title, client_logo_path):
    """Apply header logos and dynamic footer to all sections."""
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_TAB_ALIGNMENT
    from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    enfrasys_logo = SCRIPT_DIR / "enfrasys-logo.jpg"

    # --- Header ---
    doc.sections[0].different_first_page_header_footer = True
    header = doc.sections[0].header

    htable = header.add_table(1, 2, Inches(6.0))
    for p in header.paragraphs:
        if not p._element.getparent() == htable._element and p.text.strip() == "":
            p_elem_h = p._element
            p_elem_h.getparent().remove(p_elem_h)

    tblPr_h = htable._tbl.tblPr
    tblBorders = OxmlElement("w:tblBorders")
    for border_name in ["top", "left", "bottom", "right", "insideH", "insideV"]:
        border = OxmlElement(f"w:{border_name}")
        border.set(qn("w:val"), "none")
        tblBorders.append(border)
    tblPr_h.append(tblBorders)

    cell_left = htable.cell(0, 0)
    cell_right = htable.cell(0, 1)

    if client_logo_path and os.path.exists(client_logo_path):
        p_left = cell_left.paragraphs[0]
        p_left.alignment = WD_ALIGN_PARAGRAPH.LEFT
        try:
            p_left.add_run().add_picture(str(client_logo_path), height=Inches(0.6))
        except Exception:
            pass

    cell_right.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    if os.path.exists(enfrasys_logo):
        p_right = cell_right.paragraphs[0]
        p_right.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        try:
            p_right.add_run().add_picture(str(enfrasys_logo), width=Inches(1.5))
        except Exception:
            pass

    # --- Footer ---
    clean_title = document_title.replace('\n', '').replace('\r', '').strip()

    for section in doc.sections:
        footer = section.footer
        for p in footer.paragraphs:
            p_element = p._element
            p_element.getparent().remove(p_element)

        para = footer.add_paragraph()
        para.paragraph_format.tab_stops.clear_all()
        para.paragraph_format.tab_stops.add_tab_stop(Inches(6.5), WD_TAB_ALIGNMENT.RIGHT)

        run_left = para.add_run(f"{clean_title} - Planning & Design Document\t")
        run_left.font.name = 'Segoe UI'
        run_left.font.size = Pt(9)

        run_right = para.add_run()
        run_right.font.name = 'Segoe UI'
        run_right.font.size = Pt(9)

        fldChar_begin = OxmlElement('w:fldChar')
        fldChar_begin.set(qn('w:fldCharType'), 'begin')
        instrText = OxmlElement('w:instrText')
        instrText.set(qn('xml:space'), 'preserve')
        instrText.text = "PAGE"
        fldChar_end = OxmlElement('w:fldChar')
        fldChar_end.set(qn('w:fldCharType'), 'end')

        run_right._r.append(fldChar_begin)
        run_right._r.append(instrText)
        run_right._r.append(fldChar_end)


def _build_front_matter_and_toc(doc, first_p, client_name, project_title):
    """Build Change Record, Distribution List, Legal, TOC, and Sign-Off pages."""
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    anchor_p = first_p.insert_paragraph_before("")
    p_xml = anchor_p._element

    def add_p_before(text="", bold=False, size=None, space_before=None):
        p = anchor_p.insert_paragraph_before()
        if space_before is not None:
            p.paragraph_format.space_before = Pt(space_before)
        if text:
            r = p.add_run(text)
            if bold: r.bold = True
            if size: r.font.size = Pt(size)
            r.font.name = 'Segoe UI'
        return p

    # --- PAGE 2: Change Record & Distribution ---
    add_p_before("Change Record", bold=True, size=14)
    table1 = doc.add_table(rows=3, cols=4)
    try: table1.style = 'Table Grid'
    except KeyError: pass
    hdr_cells = table1.rows[0].cells
    hdr_cells[0].text = 'Date'
    hdr_cells[1].text = 'Author'
    hdr_cells[2].text = 'Version'
    hdr_cells[3].text = 'Change Reference'
    row_cells = table1.rows[1].cells
    row_cells[0].text = '[change-record-date]'
    row_cells[1].text = '[change-record-author]'
    row_cells[2].text = '[change-record-version]'
    row_cells[3].text = '[change-record-reference]'
    p_xml.addprevious(table1._element)

    add_p_before("", space_before=12)
    add_p_before("Distribution List", bold=True, size=14)
    table2 = doc.add_table(rows=3, cols=2)
    try: table2.style = 'Table Grid'
    except KeyError: pass
    hdr_cells2 = table2.rows[0].cells
    hdr_cells2[0].text = 'Name'
    hdr_cells2[1].text = 'Position/ Team'
    row_cells2 = table2.rows[1].cells
    row_cells2[0].text = '[Client/Entity/Company Name]'
    row_cells2[1].text = '[Client/Entity/Company Position]'
    row_cells3 = table2.rows[2].cells
    row_cells3[0].text = ''
    row_cells3[1].text = ''
    p_xml.addprevious(table2._element)

    add_p_before("", space_before=24)

    def add_legal(title, text):
        p = add_p_before(title, bold=True, size=7, space_before=12)
        p.paragraph_format.left_indent = Inches(2.3)
        ptext = add_p_before(text, size=7)
        ptext.paragraph_format.left_indent = Inches(2.3)
        ptext.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

    add_legal("Terms of Use", "This work may be used \"as-is\" by any interested party. You may copy, adapt, and redistribute this document for non-commercial use or for your own internal use in a commercial setting. However, you may not republish this document, nor may you publish or distribute any adaptation of this document for other than non-commercial use or your own internal use, without first obtaining express written approval from Enfrasys Solutions Sdn Bhd.")
    add_legal("Disclaimer", "The Author and Enfrasys Solutions Sdn Bhd shall have neither liability nor responsibility to any person or entity with respect to the loss or damages arising from the information contained in this work. This work may include inaccuracies or typographical errors and solely represent the opinions of the Author. Changes are periodically made to this document without notice.Due to the rapid growth of various technologies, the Author and Enfrasys Solutions Sdn Bhd cannot guarantee the accuracy of the information presented after the date of publication.")
    add_legal("Trademarks", "The names of actual companies, service marks, trademarks or products mentioned herein may be the trademarks of their respective owners. Use of terms within this work should not be regarded as affecting the validity of any trademark or service mark. Enfrasys Consulting may have patents, patent applications, trademarks, copyrights, or other intellectual property rights covering subject matter in this document. Except as expressly provided in any written license agreement from Enfrasys Solutions Sdn Bhd, the furnishing of this document does not give you any license to those items.")

    p_break2 = add_p_before("")
    p_break2.add_run().add_break(WD_BREAK.PAGE)

    # --- PAGE 3: TABLE OF CONTENTS (XML INJECTION) ---
    p_toc_title = add_p_before("Table of Contents", size=15, bold=True)
    p_toc_title.paragraph_format.space_after = Pt(12)

    p_toc = add_p_before("")
    run_toc = p_toc.add_run()

    fldChar1 = OxmlElement('w:fldChar')
    fldChar1.set(qn('w:fldCharType'), 'begin')
    instrText = OxmlElement('w:instrText')
    instrText.set(qn('xml:space'), 'preserve')
    instrText.text = 'TOC \\o "1-3" \\h \\z \\u'
    fldChar2 = OxmlElement('w:fldChar')
    fldChar2.set(qn('w:fldCharType'), 'separate')
    fldChar3 = OxmlElement('w:fldChar')
    fldChar3.set(qn('w:fldCharType'), 'end')

    run_toc._r.append(fldChar1)
    run_toc._r.append(instrText)
    run_toc._r.append(fldChar2)
    run_toc._r.append(fldChar3)

    p_break_toc = add_p_before("")
    p_break_toc.add_run().add_break(WD_BREAK.PAGE)

    # Force Word/LibreOffice to update fields on open
    doc_settings = doc.settings.element
    update_fields = OxmlElement('w:updateFields')
    update_fields.set(qn('w:val'), 'true')
    doc_settings.append(update_fields)

    # --- PAGE 4: SIGN OFF ---
    p_signoff = add_p_before("1.0 Design Document Sign Off", size=15)
    try: p_signoff.style = 'Heading 1'
    except Exception: pass

    signoff_text_1 = f"We hereby acknowledge that the design document for {client_name} {project_title} has been reviewed, and all key aspects have been addressed satisfactorily."
    add_p_before(signoff_text_1, size=11, space_before=12)
    signoff_text_2 = "This document has been prepared, reviewed, accepted, and signed off by the following individuals:"
    add_p_before(signoff_text_2, size=11, space_before=12)
    add_p_before("", space_before=4)

    def set_table_borders(table):
        tbl_borders = OxmlElement('w:tblBorders')
        for border_name in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
            border = OxmlElement(f'w:{border_name}')
            border.set(qn('w:val'), 'single')
            border.set(qn('w:sz'), '8')
            border.set(qn('w:color'), '000000')
            tbl_borders.append(border)
        table._tbl.tblPr.append(tbl_borders)

    # Signature table 1
    table_sig_1 = doc.add_table(rows=1, cols=2)
    try: table_sig_1.style = 'Table Grid'
    except Exception: pass
    set_table_borders(table_sig_1)
    s1_cells = table_sig_1.rows[0].cells
    p_s1_left = s1_cells[0].paragraphs[0]
    p_s1_left.add_run("Prepared by:\n\n")
    p_s1_left.add_run("Enfrasys Solutions Sdn Bhd").bold = True
    p_s1_left.add_run("\n\nSignature:\n\n\n\n\n___________________________\nName:\nDesignation:\nDate:")
    p_s1_right = s1_cells[1].paragraphs[0]
    p_s1_right.add_run("Verified by:\n\n")
    p_s1_right.add_run("Enfrasys Solutions Sdn Bhd").bold = True
    p_s1_right.add_run("\n\nSignature:\n\n\n\n\n___________________________\nName:\nDesignation:\nDate:")
    p_xml.addprevious(table_sig_1._element)

    p_break3 = add_p_before("")
    p_break3.add_run().add_break(WD_BREAK.PAGE)

    # Signature table 2
    table_sig_2 = doc.add_table(rows=4, cols=2)
    try: table_sig_2.style = 'Table Grid'
    except Exception: pass
    set_table_borders(table_sig_2)
    r1_cells = table_sig_2.rows[0].cells
    r1_cells[0].merge(r1_cells[1])
    r1_cells[0].text = ""
    p_r1 = r1_cells[0].paragraphs[0]
    p_r1.add_run("Reviewed by:\n").bold = True
    p_r1.add_run("\n\n")
    r2_cells = table_sig_2.rows[1].cells
    r2_cells[0].text = ""
    r2_cells[0].paragraphs[0].add_run("Name").bold = True
    r2_cells[1].text = ""
    r2_cells[1].paragraphs[0].add_run("Signature").bold = True
    r3_cells = table_sig_2.rows[2].cells
    r3_cells[0].text = "Name:\n\nDesignation:\n\nDate:"
    r3_cells[1].text = "\n\n\n\n\n___________________________"
    r4_cells = table_sig_2.rows[3].cells
    r4_cells[0].text = ""
    p_r4_L = r4_cells[0].paragraphs[0]
    p_r4_L.add_run("Verified by:\n\n").bold = True
    p_r4_L.add_run(f"[Designation]\n\n{client_name}\n\nSignature:\n\n\n\n\n___________________________\nName:\n\nDesignation:\n\nDate:")
    r4_cells[1].text = ""
    p_r4_R = r4_cells[1].paragraphs[0]
    p_r4_R.add_run("Approved by:\n\n").bold = True
    p_r4_R.add_run(f"[Designation]\n\n{client_name}\n\nSignature:\n\n\n\n\n___________________________\nName:\n\nDesignation:\n\nDate:")
    p_xml.addprevious(table_sig_2._element)

    p_break4 = add_p_before("")
    p_break4.add_run().add_break(WD_BREAK.PAGE)

    # Remove the anchor paragraph
    p_elem_anchor = anchor_p._element
    p_elem_anchor.getparent().remove(p_elem_anchor)


def _fix_tables(doc):
    """Enterprise-grade table fix: explicit XML borders, fixed layout, header banding.

    This is the key fix for the 'scattered table' problem. We explicitly set every
    border and column width via OxmlElement, bypassing style-based rendering which
    breaks when LibreOffice doesn't have the referenced Word style.
    """
    from docx.shared import Pt, Inches, Emu, RGBColor
    from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    ENFRASYS_BLUE = "4472C4"
    ENFRASYS_BLUE_LIGHT = "D6E4F0"
    PAGE_WIDTH_EMU = Inches(6.5)  # Standard page width minus margins

    for table in doc.tables:
        tblPr = table._tbl.tblPr
        if tblPr is None:
            tblPr = OxmlElement('w:tblPr')
            table._tbl.insert(0, tblPr)

        # --- Detect if this is a signature/front-matter table ---
        is_signature_table = False
        for row in table.rows:
            for cell in row.cells:
                cell_text = cell.text.lower()
                if "prepared by:" in cell_text or "reviewed by:" in cell_text or "verified by:" in cell_text:
                    is_signature_table = True
                    break
            if is_signature_table:
                break

        # --- 1. Remove any existing borders and set explicit ones ---
        existing_borders = tblPr.find(qn('w:tblBorders'))
        if existing_borders is not None:
            tblPr.remove(existing_borders)

        tblBorders = OxmlElement('w:tblBorders')
        border_color = "000000" if is_signature_table else ENFRASYS_BLUE
        for border_name in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
            border = OxmlElement(f'w:{border_name}')
            border.set(qn('w:val'), 'single')
            border.set(qn('w:sz'), '4')
            border.set(qn('w:space'), '0')
            border.set(qn('w:color'), border_color)
            tblBorders.append(border)
        tblPr.append(tblBorders)

        # --- 2. Set table width to 100% of page ---
        existing_tblW = tblPr.find(qn('w:tblW'))
        if existing_tblW is not None:
            tblPr.remove(existing_tblW)
        tblW = OxmlElement('w:tblW')
        tblW.set(qn('w:type'), 'pct')
        tblW.set(qn('w:w'), '5000')  # 5000 = 100% in Word's percentage units
        tblPr.append(tblW)

        # --- 3. Set fixed table layout (prevents column collapse in LibreOffice) ---
        existing_layout = tblPr.find(qn('w:tblLayout'))
        if existing_layout is not None:
            tblPr.remove(existing_layout)
        tblLayout = OxmlElement('w:tblLayout')
        tblLayout.set(qn('w:type'), 'fixed')
        tblPr.append(tblLayout)

        # --- 4. Distribute column widths evenly ---
        num_cols = len(table.columns)
        if num_cols > 0:
            col_width = int(PAGE_WIDTH_EMU / num_cols)
            for row in table.rows:
                for idx, cell in enumerate(row.cells):
                    tcPr = cell._tc.get_or_add_tcPr()
                    # Set cell width
                    existing_tcW = tcPr.find(qn('w:tcW'))
                    if existing_tcW is not None:
                        tcPr.remove(existing_tcW)
                    tcW = OxmlElement('w:tcW')
                    tcW.set(qn('w:type'), 'dxa')
                    tcW.set(qn('w:w'), str(int(col_width / 914)))  # Convert EMU to twips
                    tcPr.append(tcW)

        # --- 5. Apply header row banding (non-signature tables only) ---
        if not is_signature_table and len(table.rows) > 0:
            # Header row: blue background, white bold text
            for cell in table.rows[0].cells:
                tcPr = cell._tc.get_or_add_tcPr()
                # Remove existing shading
                existing_shd = tcPr.find(qn('w:shd'))
                if existing_shd is not None:
                    tcPr.remove(existing_shd)
                shading = OxmlElement('w:shd')
                shading.set(qn('w:val'), 'clear')
                shading.set(qn('w:color'), 'auto')
                shading.set(qn('w:fill'), ENFRASYS_BLUE)
                tcPr.append(shading)

                for p in cell.paragraphs:
                    for run in p.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.bold = True
                        run.font.name = 'Segoe UI'
                        run.font.size = Pt(10)

            # Alternating row banding (light blue for even data rows)
            for row_idx, row in enumerate(table.rows):
                if row_idx == 0:
                    continue  # Skip header
                if row_idx % 2 == 0:  # Even rows get light shading
                    for cell in row.cells:
                        tcPr = cell._tc.get_or_add_tcPr()
                        existing_shd = tcPr.find(qn('w:shd'))
                        if existing_shd is not None:
                            tcPr.remove(existing_shd)
                        shading = OxmlElement('w:shd')
                        shading.set(qn('w:val'), 'clear')
                        shading.set(qn('w:color'), 'auto')
                        shading.set(qn('w:fill'), ENFRASYS_BLUE_LIGHT)
                        tcPr.append(shading)

            # Set header row to repeat on page break
            first_tr = table.rows[0]._tr
            trPr = first_tr.find(qn('w:trPr'))
            if trPr is None:
                trPr = OxmlElement('w:trPr')
                first_tr.insert(0, trPr)
            tblHeader = OxmlElement('w:tblHeader')
            trPr.append(tblHeader)

        # --- 6. Cell formatting: vertical alignment, padding, fonts ---
        for row in table.rows:
            for cell in row.cells:
                cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER

                # Set cell margins/padding via XML
                tcPr = cell._tc.get_or_add_tcPr()
                tcMar = OxmlElement('w:tcMar')
                for margin_name in ['top', 'bottom', 'left', 'right']:
                    margin = OxmlElement(f'w:{margin_name}')
                    margin.set(qn('w:w'), '40')  # ~0.03 inches
                    margin.set(qn('w:type'), 'dxa')
                    tcMar.append(margin)
                # Remove existing margins first
                existing_mar = tcPr.find(qn('w:tcMar'))
                if existing_mar is not None:
                    tcPr.remove(existing_mar)
                tcPr.append(tcMar)

                for p in cell.paragraphs:
                    p.paragraph_format.space_before = Pt(1)
                    p.paragraph_format.space_after = Pt(1)
                    for run in p.runs:
                        run.font.name = 'Segoe UI'
                        if run.font.size is None:
                            run.font.size = Pt(10)


def _fix_typography_and_captions(doc):
    """Enforce global font, heading sizes, caption styling, page breaks, and text justification."""
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    style_rules = {
        'Heading 1': {'size': 15, 'bold': False},
        'Heading 2': {'size': 14, 'bold': False},
        'Heading 3': {'size': 13, 'bold': False},
        'Heading 4': {'size': 12, 'bold': False},
    }

    figure_counter = 1

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # --- A. FIX IMAGE CAPTIONS (Renumber & Center Safely) ---
        if text.lower().startswith("figure ") and ":" in text:
            new_text = re.sub(r'(?i)^figure\s+\d+\s*:', f'Figure {figure_counter}:', text)

            # SAFE TEXT REPLACEMENT: Protect inline images from being deleted!
            has_drawing = any(getattr(run._element, "drawing_lst", None) for run in para.runs)

            if has_drawing:
                for run in para.runs:
                    if run.text:
                        run.text = ""
                para.add_run("\n" + new_text)
            else:
                para.text = new_text

            para.alignment = WD_ALIGN_PARAGRAPH.CENTER

            for run in para.runs:
                if run.text.strip():
                    run.font.name = 'Segoe UI'
                    run.font.size = Pt(10)
                    run.font.italic = True
                    run.font.color.rgb = RGBColor(128, 128, 128)

            figure_counter += 1
            continue

        # --- B. STANDARD TEXT JUSTIFICATION ---
        if para.style.name == "Normal":
            # Skip justification if the paragraph contains an image to preserve image alignments
            has_drawing = any(getattr(run._element, "drawing_lst", None) for run in para.runs)
            if not has_drawing:
                para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

        # --- C. PAGE BREAKS FOR MAIN SECTIONS ---
        if para.style.name == "Heading 1" or any(text.startswith(f"{i}.0") for i in range(2, 21)):
            if not text.startswith("1.0") and not text.startswith("2.0"):
                para.paragraph_format.page_break_before = True

        # --- D. ENFORCE TYPOGRAPHY SCALING ---
        style_name = para.style.name
        rule = style_rules.get(style_name)
        for run in para.runs:
            run.font.name = 'Segoe UI'
            if rule:
                run.font.size = Pt(rule['size'])
                run.font.bold = rule['bold']
            else:
                if run.font.size is None:
                    run.font.size = Pt(11)


def _fix_images(doc):
    """Center images and enforce max height with proportional scaling."""
    from docx.shared import Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    # Center image paragraphs (unless explicitly left aligned)
    for p in doc.paragraphs:
        for run in p.runs:
            drawing = getattr(run._element, "drawing_lst", None)
            if drawing:
                if p.alignment != WD_ALIGN_PARAGRAPH.LEFT:
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Resize oversized images
    for shape in doc.inline_shapes:
        max_height = Inches(4.5)
        if shape.height > max_height:
            ratio = max_height / float(shape.height)
            shape.height = max_height
            shape.width = int(shape.width * ratio)


def convert_md_to_docx(md_path, docx_path, document_title, project_title, client_name, client_logo_path=None):
    """Master function: Markdown → DOCX with full Enfrasys formatting.
    
    Pipeline:
    1. Pandoc converts MD → raw DOCX (with template styles)
    2. python-docx post-processes: title page, headers, footers, front matter, TOC, sign-off
    3. Enterprise table fixes (explicit XML borders, fixed layout, header banding)
    4. Typography enforcement (fonts, captions, page breaks)
    5. LibreOffice headless updates TOC fields
    """
    try:
        from docx import Document

        # Phase 1: Pandoc conversion
        _pandoc_md_to_raw_docx(str(md_path), str(docx_path))

        # Phase 2: Load and post-process
        doc = Document(str(docx_path))

        # Prepare first paragraph for front-page insertion
        if len(doc.paragraphs) > 0:
            first_p = doc.paragraphs[0]
            p_element = first_p._p
            for child in list(p_element):
                p_element.remove(child)
        else:
            first_p = doc.add_paragraph()

        # Build document structure
        _build_title_page(doc, first_p, project_title, document_title, client_name, client_logo_path)
        _apply_header_footer(doc, document_title, client_logo_path)
        _build_front_matter_and_toc(doc, first_p, client_name, project_title)

        # Phase 3: Fix content formatting
        _fix_images(doc)
        _fix_tables(doc)
        _fix_typography_and_captions(doc)

        # Save intermediate DOCX
        doc.save(str(docx_path))

        # Phase 4: Update TOC fields via LibreOffice
        try:
            _update_toc_with_libreoffice(str(docx_path))
            print("[convert_md_to_docx] TOC updated via LibreOffice roundtrip.")
        except Exception as e:
            print(f"[convert_md_to_docx] TOC roundtrip failed/skipped: {e}")

    except Exception as e:
        print(f"[convert_md_to_docx] Warning: {e}")
        import traceback
        traceback.print_exc()


# -----------------------------
# API Endpoints
# -----------------------------
@app.post("/api/upload")
async def upload_file(background_tasks: BackgroundTasks, file: UploadFile = File(...), logo: UploadFile = File(None)):
    task_id = str(uuid.uuid4())
    task_dir = UPLOAD_DIR / task_id
    task_dir.mkdir(parents=True, exist_ok=True)

    src_name = safe_filename(file.filename, default="source.pptx")
    file_path = task_dir / f"source_{src_name}"

    with open(file_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    logo_path = None
    if logo:
        logo_name = safe_filename(logo.filename, default="logo")
        logo_path = task_dir / f"logo_{logo_name}"
        with open(logo_path, "wb") as f:
            shutil.copyfileobj(logo.file, f)

    # Register task in DB
    update_task(task_id, {
        "status": "processing_upload",
        "step_name": "Saving files to Azure...",
        "progress": 5,
        "cost_metrics": {
            "vision_tokens_prompt": 0,
            "vision_tokens_completion": 0,
            "llm_tokens_prompt": 0,
            "llm_tokens_completion": 0,
            "content_understanding_pages": 0,
            "total_cost_myr": 0.0
        }
    })

    # Trigger the background worker
    background_tasks.add_task(
        process_initial_upload,
        task_id,
        str(file_path),
        str(logo_path) if logo_path else None,
        src_name,
        task_dir
    )

    # Return INSTANTLY to Next.js before Vercel times out!
    return {"task_id": task_id, "status": "processing_upload"}


def background_processing(task_id: str):
    task = get_task(task_id)
    if not task:
        return

    file_path = task["file_path"]
    logo_path = task.get("logo_path")
    task_dir = UPLOAD_DIR / task_id

    update_task(task_id, {
        "status": "processing",
        "step_name": "Extracting Structure & Images",
        "progress": 15
    })

    try:
        working_pptx_path = str(task_dir / "working.pptx")
        working_pdf_path = str(task_dir / "working.pdf")

        shutil.copy2(file_path, working_pptx_path)
        pptx_data = extract_with_pptx(working_pptx_path, task_dir)

        update_task(task_id, {
            "step_name": "Converting PPTX to PDF & AI Analysis",
            "progress": 30
        })

        convert_pptx_to_pdf(working_pptx_path, working_pdf_path)
        cu_data = analyze_with_content_understanding(working_pdf_path)

        update_task(task_id, {
            "step_name": "Analyzing Images with GPT-4o Vision",
            "progress": 50
        })

        pptx_data["slides"] = analyze_images_with_vision(task_id, pptx_data["slides"])

        merged = merge_extraction_results(pptx_data, cu_data)
        payload_file = task_dir / "extraction_payload.json"
        with open(payload_file, "w", encoding="utf-8") as f:
            json.dump(merged, f, indent=2, ensure_ascii=False, default=str)

        update_task(task_id, {
            "step_name": "Orchestrating Document Structure",
            "progress": 70
        })

        toc = generate_table_of_contents(task_id, merged)
        if not toc:
            raise Exception("Failed to generate Table of Contents.")

        update_task(task_id, {
            "step_name": "Drafting Document Sections (Taking 2-3 mins)",
            "progress": 85
        })

        final_doc_md = write_document_sections(task_id, toc, merged)
        final_doc_md = final_doc_md.replace(str(task_dir) + os.sep, "")

        md_file = task_dir / "FINAL_DESIGN_DOCUMENT.md"
        with open(md_file, "w", encoding="utf-8") as f:
            f.write(final_doc_md)

        update_task(task_id, {
            "step_name": "Converting to Word Document",
            "progress": 95
        })

        docx_path = task_dir / "Solution_Design_Document.docx"
        client_name = toc.get("client_name", "").strip()
        project_title = toc.get("project_title", "Project").strip()
        doc_title = f"{client_name} {project_title}" if client_name else project_title

        convert_md_to_docx(md_file, docx_path, doc_title, project_title, client_name, logo_path)

        asset_library = []
        for slide in merged.get("slides", []):
            for img in slide.get("images", []):
                if img.get("ai_description") and img.get("ai_description") != "DECORATIVE":
                    asset_library.append({
                        "filename": img.get("filename"),
                        "description": img.get("ai_description")
                    })

        markdown_text = ""
        if md_file.exists():
            with open(md_file, "r", encoding="utf-8") as f:
                markdown_text = f.read()

        # Final Cost Calculation
        final_task = get_task(task_id)
        if final_task and "cost_metrics" in final_task:
            metrics = final_task["cost_metrics"]
            # Track content understanding pages based on slide count
            metrics["content_understanding_pages"] = len(merged.get("slides", []))
            
            # Calculate USD values
            cost_usd = 0.0
            cost_usd += metrics["vision_tokens_prompt"] * RATE_VISION_PROMPT
            cost_usd += metrics["vision_tokens_completion"] * RATE_VISION_COMPLETION
            cost_usd += metrics["llm_tokens_prompt"] * RATE_LLM_PROMPT
            cost_usd += metrics["llm_tokens_completion"] * RATE_LLM_COMPLETION
            cost_usd += metrics["content_understanding_pages"] * RATE_CU_PER_PAGE
            
            metrics["total_cost_myr"] = round(cost_usd * USD_TO_MYR_RATE, 2)
            
            update_task(task_id, {"cost_metrics": metrics})

        update_task(task_id, {
            "status": "completed",
            "step_name": "Ready!",
            "progress": 100,
            "result_docx": str(docx_path.resolve()),
            "markdown_draft": markdown_text,
            "asset_library": asset_library
        })

    except Exception as e:
        update_task(task_id, {
            "status": "failed",
            "step_name": "Failed",
            "progress": 0,
            "error": str(e)
        })


@app.get("/api/thumbnails/{task_id}/{filename}")
async def get_thumbnail(task_id: str, filename: str):
    file_path = UPLOAD_DIR / task_id / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="Thumbnail not found")
    return FileResponse(str(file_path))


@app.post("/api/process/{task_id}")
async def start_processing(task_id: str, background_tasks: BackgroundTasks):
    if not get_task(task_id):
        return {"error": "Invalid task ID"}
    background_tasks.add_task(background_processing, task_id)
    return {"status": "started", "task_id": task_id}


@app.get("/api/status/{task_id}")
async def get_status(task_id: str):
    task = get_task(task_id)
    if not task:
        return {"error": "Invalid task ID"}

    response = {
        "status": task.get("status"),
        "step_name": task.get("step_name"),
        "progress": task.get("progress")
    }

    if task.get("status") == "completed":
        response["markdown_draft"] = task.get("markdown_draft")
        response["asset_library"] = task.get("asset_library", [])
        response["page_count"] = task.get("page_count", 0)

    # Add this inside get_status
    if task.get("status") == "upload_complete":
        response["preview_data"] = task.get("preview_data")

    if "cost_metrics" in task:
        response["cost_metrics"] = task["cost_metrics"]

    return response


@app.get("/api/download/{task_id}")
async def download_doc(task_id: str, filename: str = "Enfrasys_Design_Document"):
    task = get_task(task_id)
    if not task:
        return {"error": "Invalid task ID"}

    if task.get("status") != "completed" or "result_docx" not in task:
        return {"error": "Document not ready or generation failed."}

    return FileResponse(
        path=task["result_docx"],
        filename=f"{filename}.docx",
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )


@app.get("/api/download-pdf/{task_id}")
def download_pdf(task_id: str, filename: str = "Enfrasys_Design_Document"):
    task = get_task(task_id)
    if not task:
        raise HTTPException(status_code=404, detail="Invalid task ID")

    if task.get("status") != "completed":
        raise HTTPException(status_code=400, detail="Document not ready.")

    if task.get("result_pdf") and Path(task["result_pdf"]).exists():
        return FileResponse(
            path=task["result_pdf"],
            filename=f"{filename}.pdf",
            media_type="application/pdf"
        )

    docx_path = Path(task["result_docx"])
    pdf_path = docx_path.parent / "Solution_Design_Document.pdf"

    try:
        convert_docx_to_pdf(str(docx_path.resolve()), str(pdf_path.resolve()))
        if not pdf_path.exists():
            raise Exception("PDF file was not created by the converter.")
        update_task(task_id, {"result_pdf": str(pdf_path.resolve())})
        task["result_pdf"] = str(pdf_path.resolve())
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF generation failed: {str(e)}")

    return FileResponse(
        path=task["result_pdf"],
        filename=f"{filename}.pdf",
        media_type="application/pdf"
    )


@app.get("/api/prepare-preview/{task_id}")
def prepare_preview(task_id: str):
    task = get_task(task_id)
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")

    if task.get("status") != "completed":
        raise HTTPException(status_code=400, detail="Task processing not completed yet")

    if task.get("preview_prepared"):
        return {
            "page_count": task.get("page_count", 0),
            "page_images": task.get("page_images", []),
            "cost_metrics": task.get("cost_metrics", {})
        }

    if "result_docx" not in task:
        raise HTTPException(status_code=404, detail="Result DOCX path not found in task")

    docx_path = Path(task.get("result_docx"))
    if not docx_path.exists():
        raise HTTPException(status_code=404, detail="Generated DOCX not found")

    task_dir = UPLOAD_DIR / task_id
    pdf_path = task_dir / "Solution_Design_Document.pdf"

    try:
        convert_docx_to_pdf(str(docx_path.resolve()), str(pdf_path.resolve()))
        if not pdf_path.exists():
            raise Exception("PDF file was not created by the converter.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Word to PDF conversion failed: {str(e)}")

    pages_dir = task_dir / "doc_pages"
    pages_dir.mkdir(parents=True, exist_ok=True)

    page_image_urls = []
    page_count = 0

    try:
        import fitz
        pdf_doc = fitz.open(str(pdf_path))
        page_count = len(pdf_doc)
        for i in range(page_count):
            page = pdf_doc[i]
            pix = page.get_pixmap(dpi=150)
            img_path = pages_dir / f"page_{i + 1}.jpg"
            pix.save(str(img_path))
            page_image_urls.append(f"/api/doc-pages/{task_id}/page_{i + 1}.jpg")
        pdf_doc.close()
    except ImportError:
        raise HTTPException(status_code=500, detail="PyMuPDF not installed. pip install PyMuPDF")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF to image conversion failed: {str(e)}")

    update_task(task_id, {
        "page_count": page_count,
        "page_images": page_image_urls,
        "result_pdf": str(pdf_path.resolve()),
        "preview_prepared": True
    })

    return {
        "page_count": page_count,
        "page_images": page_image_urls,
        "cost_metrics": task.get("cost_metrics", {})
    }


@app.get("/api/doc-pages/{task_id}/{filename}")
async def get_doc_page(task_id: str, filename: str):
    file_path = UPLOAD_DIR / task_id / "doc_pages" / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="Page image not found")
    return FileResponse(str(file_path), media_type="image/jpeg")


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=int(os.getenv("PORT", "8000")))