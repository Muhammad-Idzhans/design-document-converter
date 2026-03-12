"""
Phase 1: Hybrid Extraction Pipeline
=====================================
Extracts content from PPTX design documents using:
1. python-pptx  — structured extraction of text, tables, speaker notes, and images
2. Azure Content Understanding — AI-powered analysis of slides (including diagram understanding)

Outputs a structured JSON payload ready for Phase 2 (AI Generation).

Usage:
    python slides-to-doc.py <path-to-pptx>
    python slides-to-doc.py  (defaults to sample document)
"""

import os
import sys
import json
import time
import base64
import requests
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from dotenv import load_dotenv

# ──────────────────────────────────────────────
# Configuration
# ──────────────────────────────────────────────
load_dotenv()

CONTENT_UNDERSTANDING_ENDPOINT = os.getenv("CONTENT_UNDERSTANDING_ENDPOINT", "").rstrip("/")
CONTENT_UNDERSTANDING_KEY = os.getenv("CONTENT_UNDERSTANDING_KEY", "")
API_VERSION = "2025-11-01"

SCRIPT_DIR = Path(__file__).parent
PROJECT_DIR = SCRIPT_DIR.parent
DEFAULT_PPTX = PROJECT_DIR / "sample-documents" / "UMS_PowerBI_Fabric_Design_Workshop_v1.1_final.pptx"
OUTPUT_DIR = PROJECT_DIR / "extraction-output"


# ──────────────────────────────────────────────
# Part 1: python-pptx Extraction
# ──────────────────────────────────────────────
def extract_with_pptx(pptx_path: str) -> dict:
    """
    Extract structured content from PPTX using python-pptx.
    Returns text, tables, speaker notes, and saves images to disk.
    """
    prs = Presentation(pptx_path)
    images_dir = OUTPUT_DIR / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    slides_data = []
    image_counter = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_info = {
            "slide_number": slide_idx,
            "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
            "title": None,
            "text_content": [],
            "speaker_notes": None,
            "images": [],
            "tables": [],
            "is_section_divider": False,
        }

        # ── Title ──
        if slide.shapes.title:
            slide_info["title"] = slide.shapes.title.text.strip()

        # ── Shapes: text, images, tables ──
        has_body_text = False
        for shape in slide.shapes:
            # Text
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text != slide_info["title"]:
                        slide_info["text_content"].append(text)
                        has_body_text = True

            # Images
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_counter += 1
                img = shape.image
                ext = img.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                filename = f"slide_{slide_idx:03d}_img_{image_counter:03d}.{ext}"
                filepath = images_dir / filename

                with open(filepath, "wb") as f:
                    f.write(img.blob)

                slide_info["images"].append({
                    "image_id": f"img_{slide_idx:03d}_{image_counter:03d}",
                    "filename": filename,
                    "file_path": str(filepath),
                    "content_type": img.content_type,
                    "size_bytes": len(img.blob),
                    "width_emu": shape.width if hasattr(shape, 'width') else None,
                    "height_emu": shape.height if hasattr(shape, 'height') else None,
                })

            # Tables
            if shape.has_table:
                table = shape.table
                headers = [cell.text.strip() for cell in table.rows[0].cells]
                rows = []
                for row_idx, row in enumerate(table.rows):
                    if row_idx == 0:
                        continue  # skip header
                    row_data = [cell.text.strip() for cell in row.cells]
                    rows.append(row_data)

                slide_info["tables"].append({
                    "headers": headers,
                    "rows": rows,
                    "total_rows": len(table.rows),
                    "total_cols": len(table.columns),
                })

        # ── Speaker Notes ──
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_info["speaker_notes"] = notes

        # ── Detect section dividers ──
        # Section divider slides typically have a title but minimal body content
        if slide_info["title"] and not has_body_text and not slide_info["tables"] and not slide_info["images"]:
            slide_info["is_section_divider"] = True

        slides_data.append(slide_info)

    # Build metadata
    pptx_filename = Path(pptx_path).name
    project_name = None
    if slides_data and slides_data[0]["title"]:
        project_name = slides_data[0]["title"]

    return {
        "metadata": {
            "source_file": pptx_filename,
            "total_slides": len(prs.slides),
            "total_images_extracted": image_counter,
            "total_tables_extracted": sum(len(s["tables"]) for s in slides_data),
            "slides_with_notes": sum(1 for s in slides_data if s["speaker_notes"]),
            "project_name": project_name,
        },
        "slides": slides_data,
    }


# ──────────────────────────────────────────────
# Part 2: Azure Content Understanding Analysis
# ──────────────────────────────────────────────
def analyze_with_content_understanding(pptx_path: str) -> dict | None:
    """
    Send the PPTX to Azure Content Understanding for AI-powered analysis.
    Uses the design_document_converter analyzer via the REST API.
    Returns the full analysis result including markdown content and extracted fields.
    """
    if not CONTENT_UNDERSTANDING_ENDPOINT or not CONTENT_UNDERSTANDING_KEY:
        print("[INFO] Content Understanding credentials not configured. Skipping AI analysis.")
        print("[INFO] Set CONTENT_UNDERSTANDING_ENDPOINT and CONTENT_UNDERSTANDING_KEY in .env")
        return None

    print(f"[CU] Sending PPTX to Content Understanding for analysis...")

    # ── Step 1: Submit the file for analysis (binary upload) ──
    analyze_url = (
        f"{CONTENT_UNDERSTANDING_ENDPOINT}/contentunderstanding"
        f"/analyzers/design_document_converter:analyze?api-version={API_VERSION}"
    )

    with open(pptx_path, "rb") as f:
        file_bytes = f.read()

    # Use analyzeBinary for direct file upload
    analyze_binary_url = (
        f"{CONTENT_UNDERSTANDING_ENDPOINT}/contentunderstanding"
        f"/analyzers/design_document_converter:analyzeBinary?api-version={API_VERSION}"
    )

    headers = {
        "Ocp-Apim-Subscription-Key": CONTENT_UNDERSTANDING_KEY,
        "Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }

    response = requests.post(analyze_binary_url, headers=headers, data=file_bytes)

    if response.status_code not in (200, 202):
        print(f"[CU] Error submitting file: {response.status_code}")
        print(f"[CU] Response: {response.text}")

        # Fallback: try URL-based method if binary fails
        print(f"[CU] Trying base64 URL method as fallback...")
        return _analyze_with_base64_fallback(pptx_path)

    # ── Step 2: Poll for results ──
    operation_url = response.headers.get("Operation-Location")
    if not operation_url:
        print(f"[CU] No Operation-Location in response headers.")
        print(f"[CU] Headers: {dict(response.headers)}")
        # If the result is immediate
        if response.status_code == 200:
            return response.json()
        return None

    print(f"[CU] Analysis submitted. Polling for results...")
    return _poll_for_result(operation_url)


def _analyze_with_base64_fallback(pptx_path: str) -> dict | None:
    """Fallback: send file as base64 data URI."""
    analyze_url = (
        f"{CONTENT_UNDERSTANDING_ENDPOINT}/contentunderstanding"
        f"/analyzers/design_document_converter:analyze?api-version={API_VERSION}"
    )

    with open(pptx_path, "rb") as f:
        file_bytes = f.read()

    b64_content = base64.b64encode(file_bytes).decode("utf-8")
    data_uri = f"data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_content}"

    headers = {
        "Ocp-Apim-Subscription-Key": CONTENT_UNDERSTANDING_KEY,
        "Content-Type": "application/json",
    }

    payload = {
        "inputs": [{"data": data_uri}]
    }

    response = requests.post(analyze_url, headers=headers, json=payload)

    if response.status_code not in (200, 202):
        print(f"[CU] Fallback also failed: {response.status_code}")
        print(f"[CU] Response: {response.text}")
        return None

    operation_url = response.headers.get("Operation-Location")
    if not operation_url:
        if response.status_code == 200:
            return response.json()
        return None

    return _poll_for_result(operation_url)


def _poll_for_result(operation_url: str, max_retries: int = 60, interval: int = 3) -> dict | None:
    """Poll the Operation-Location URL until the analysis completes."""
    headers = {
        "Ocp-Apim-Subscription-Key": CONTENT_UNDERSTANDING_KEY,
    }

    for attempt in range(max_retries):
        time.sleep(interval)
        response = requests.get(operation_url, headers=headers)

        if response.status_code != 200:
            print(f"[CU] Poll error: {response.status_code} - {response.text}")
            continue

        result = response.json()
        status = result.get("status", "")

        if status == "Succeeded":
            print(f"[CU] Analysis completed successfully!")
            return result
        elif status in ("Failed", "Cancelled"):
            print(f"[CU] Analysis {status}: {result.get('error', 'Unknown error')}")
            return None
        else:
            print(f"[CU] Status: {status} (attempt {attempt + 1}/{max_retries})")

    print(f"[CU] Timed out waiting for analysis results.")
    return None


def analyze_images_with_content_understanding(slides_data: list) -> list:
    """
    Send each extracted image to Content Understanding for AI description.
    This lets CU read text inside architecture diagrams and describe what's in each picture.
    Returns the slides_data with image descriptions added.
    """
    if not CONTENT_UNDERSTANDING_ENDPOINT or not CONTENT_UNDERSTANDING_KEY:
        return slides_data

    # Collect all images across all slides
    images_to_analyze = []
    for slide in slides_data:
        for img in slide.get("images", []):
            images_to_analyze.append((slide["slide_number"], img))

    if not images_to_analyze:
        return slides_data

    print(f"[CU] Analyzing {len(images_to_analyze)} images for diagram descriptions...")

    for slide_num, img_info in images_to_analyze:
        filepath = img_info.get("file_path")
        if not filepath or not os.path.exists(filepath):
            continue

        print(f"  → Analyzing {img_info['filename']} (slide {slide_num})...")

        # Read image and send as binary
        with open(filepath, "rb") as f:
            img_bytes = f.read()

        content_type = img_info.get("content_type", "image/png")

        analyze_url = (
            f"{CONTENT_UNDERSTANDING_ENDPOINT}/contentunderstanding"
            f"/analyzers/design_document_converter:analyzeBinary?api-version={API_VERSION}"
        )

        headers = {
            "Ocp-Apim-Subscription-Key": CONTENT_UNDERSTANDING_KEY,
            "Content-Type": content_type,
        }

        try:
            response = requests.post(analyze_url, headers=headers, data=img_bytes)

            if response.status_code not in (200, 202):
                print(f"    ⚠ Failed: {response.status_code}")
                continue

            # Poll for result
            operation_url = response.headers.get("Operation-Location")
            if operation_url:
                result = _poll_for_result(operation_url, max_retries=30, interval=2)
            elif response.status_code == 200:
                result = response.json()
            else:
                continue

            if result:
                cu_result = result.get("result", result)
                contents = cu_result.get("contents", [])
                if contents:
                    markdown = contents[0].get("markdown", "")
                    img_info["ai_description"] = markdown
                    print(f"    ✓ Got description ({len(markdown)} chars)")
                else:
                    print(f"    ⚠ No content returned")

        except Exception as e:
            print(f"    ⚠ Error: {e}")

    return slides_data


# ──────────────────────────────────────────────
# Part 3: Merge Results
# ──────────────────────────────────────────────
def merge_extraction_results(pptx_data: dict, cu_data: dict | None) -> dict:
    """
    Merge python-pptx structured data with Content Understanding AI analysis.
    The CU data enriches each slide with markdown content and AI-extracted insights.
    """
    result = {
        "metadata": pptx_data["metadata"],
        "extraction_sources": ["python-pptx"],
        "slides": pptx_data["slides"],
        "content_understanding": None,
    }

    if cu_data:
        result["extraction_sources"].append("content-understanding")

        # Extract the CU analysis result
        cu_result = cu_data.get("result", cu_data)

        # Content Understanding returns rich markdown + structured data per content
        contents = cu_result.get("contents", [])
        if contents:
            cu_content = contents[0]  # First content block

            # Markdown is the PRIMARY output — contains full slide text
            # including text extracted from inside images/diagrams
            markdown_content = cu_content.get("markdown", "")

            result["content_understanding"] = {
                # Primary: the full markdown rendering of the PPTX
                "markdown": markdown_content,
                "markdown_length": len(markdown_content),
                # Secondary: structured data
                "fields": cu_content.get("fields", {}),
                "tables": cu_content.get("tables", []),
                "pages": [
                    {
                        "page_number": p.get("pageNumber"),
                        "width": p.get("width"),
                        "height": p.get("height"),
                    }
                    for p in cu_content.get("pages", [])
                ],
                "kind": cu_content.get("kind", ""),
            }

        # Token usage info
        usage = cu_data.get("usage", {})
        if usage:
            result["content_understanding"]["usage"] = usage

    return result


# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────
def main():
    # Determine input file
    if len(sys.argv) > 1:
        pptx_path = sys.argv[1]
    else:
        # Open file picker dialog
        import tkinter as tk
        from tkinter import filedialog

        root = tk.Tk()
        root.withdraw()  # Hide the root window
        root.attributes("-topmost", True)  # Bring dialog to front

        pptx_path = filedialog.askopenfilename(
            title="Select a PowerPoint file",
            initialdir=str(PROJECT_DIR / "sample-documents"),
            filetypes=[
                ("PowerPoint files", "*.pptx"),
                ("All files", "*.*"),
            ],
        )
        root.destroy()

        if not pptx_path:
            print("No file selected. Exiting.")
            sys.exit(0)

    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    print(f"=" * 60)
    print(f"  Phase 1: PPTX Content Extraction")
    print(f"  Input: {Path(pptx_path).name}")
    print(f"=" * 60)

    # Ensure output directory exists
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # ── Handle OneDrive cloud-only files ──
    # OneDrive files with ReparsePoint attribute are placeholders;
    # copying them forces a download so python-pptx can read the bytes.
    import shutil
    import tempfile

    original_pptx_path = pptx_path
    temp_dir = None
    try:
        # Try reading first few bytes to check if the file is accessible
        with open(pptx_path, "rb") as f:
            f.read(4)
    except (OSError, PermissionError):
        pass  # File is accessible or will error later

    # Always copy to temp to avoid OneDrive locking issues
    temp_dir = tempfile.mkdtemp(prefix="pptx_extract_")
    temp_pptx = os.path.join(temp_dir, Path(pptx_path).name)
    print(f"\n[0/3] Copying file to temp directory (avoids OneDrive sync issues)...")
    shutil.copy2(pptx_path, temp_pptx)
    pptx_path = temp_pptx
    print(f"  ✓ Copied to {temp_pptx}")

    # ── Step 1: python-pptx extraction ──
    print(f"\n[1/3] Extracting with python-pptx...")
    pptx_data = extract_with_pptx(pptx_path)
    meta = pptx_data["metadata"]
    print(f"  ✓ {meta['total_slides']} slides extracted")
    print(f"  ✓ {meta['total_images_extracted']} images saved to extraction-output/images/")
    print(f"  ✓ {meta['total_tables_extracted']} tables extracted")
    print(f"  ✓ {meta['slides_with_notes']} slides with speaker notes")

    # Save python-pptx results
    pptx_json_path = OUTPUT_DIR / "pptx_extraction.json"
    with open(pptx_json_path, "w", encoding="utf-8") as f:
        json.dump(pptx_data, f, indent=2, ensure_ascii=False, default=str)
    print(f"  ✓ Saved to {pptx_json_path}")

    # ── Step 2: Content Understanding analysis (full PPTX text) ──
    print(f"\n[2/4] Analyzing PPTX text with Content Understanding...")
    cu_data = analyze_with_content_understanding(pptx_path)

    if cu_data:
        # Save full JSON result
        cu_json_path = OUTPUT_DIR / "content_understanding_result.json"
        with open(cu_json_path, "w", encoding="utf-8") as f:
            json.dump(cu_data, f, indent=2, ensure_ascii=False, default=str)
        print(f"  ✓ Full result saved to {cu_json_path}")

        # Save markdown output separately — this is the key output
        cu_result = cu_data.get("result", cu_data)
        contents = cu_result.get("contents", [])
        if contents:
            markdown_content = contents[0].get("markdown", "")
            if markdown_content:
                md_path = OUTPUT_DIR / "content_understanding_output.md"
                with open(md_path, "w", encoding="utf-8") as f:
                    f.write(markdown_content)
                print(f"  ✓ Markdown output saved to {md_path}")
                print(f"  ✓ Markdown length: {len(markdown_content):,} characters")
    else:
        print(f"  ⚠ Content Understanding analysis was skipped or failed")

    # ── Step 3: Content Understanding image analysis ──
    print(f"\n[3/4] Analyzing extracted images with Content Understanding...")
    pptx_data["slides"] = analyze_images_with_content_understanding(pptx_data["slides"])

    # Re-save pptx_extraction with image descriptions added
    with open(pptx_json_path, "w", encoding="utf-8") as f:
        json.dump(pptx_data, f, indent=2, ensure_ascii=False, default=str)

    # ── Step 4: Merge results ──
    print(f"\n[4/4] Merging extraction results...")
    merged = merge_extraction_results(pptx_data, cu_data)

    merged_json_path = OUTPUT_DIR / "extraction_payload.json"
    with open(merged_json_path, "w", encoding="utf-8") as f:
        json.dump(merged, f, indent=2, ensure_ascii=False, default=str)
    print(f"  ✓ Final payload saved to {merged_json_path}")

    # ── Cleanup temp directory ──
    if temp_dir and os.path.exists(temp_dir):
        shutil.rmtree(temp_dir, ignore_errors=True)

    # ── Summary ──
    print(f"\n{'=' * 60}")
    print(f"  Extraction Complete!")
    print(f"  Sources: {', '.join(merged['extraction_sources'])}")
    print(f"  Output directory: {OUTPUT_DIR}")
    print(f"{'=' * 60}")

    # Print slide-by-slide summary
    print(f"\n  Slide Summary:")
    for slide in merged["slides"]:
        markers = []
        if slide["images"]:
            markers.append(f"{len(slide['images'])} img")
        if slide["tables"]:
            markers.append(f"{len(slide['tables'])} tbl")
        if slide["speaker_notes"]:
            markers.append("notes")
        if slide["is_section_divider"]:
            markers.append("divider")

        marker_str = f" [{', '.join(markers)}]" if markers else ""
        title = slide["title"] or "(no title)"
        print(f"    Slide {slide['slide_number']:2d}: {title}{marker_str}")

    return merged


if __name__ == "__main__":
    main()
